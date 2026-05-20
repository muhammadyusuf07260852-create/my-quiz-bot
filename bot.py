import os
import telebot
from telebot.types import ReplyKeyboardMarkup, KeyboardButton, KeyboardButtonPollType, ReplyKeyboardRemove, InlineKeyboardMarkup, InlineKeyboardButton
from dotenv import load_dotenv
import sqlite3
import database
import json
import threading
import random
import docx
import pypdf
from pptx import Presentation
import google.generativeai as genai

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                text_runs.append(run.text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                text_runs.append(cell.text)
        return "\n".join(text_runs).strip()
    except Exception as e:
        print(f"PPTX extract error: {e}")
        return ""

load_dotenv()
BOT_TOKEN = os.getenv('BOT_TOKEN')
bot = telebot.TeleBot(BOT_TOKEN)

# Gemini AI konfiguratsiyasi
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# AI vaqtinchalik sessiyalari
ai_sessions = {}

def extract_text_from_pdf(file_path):
    try:
        reader = pypdf.PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"PDF extract error: {e}")
        return ""

def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        text = ""
        for para in doc.paragraphs:
            if para.text:
                text += para.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Docx extract error: {e}")
        return ""

bot_info = bot.get_me()

# Telegram botda buyruqlar menyusini o'rnatish (/ bosganda chiqadi)
bot.set_my_commands([
    telebot.types.BotCommand('/start',        '🏠 Bosh menyu'),
    telebot.types.BotCommand('/newquiz',      '➕ Yangi test yaratish'),
    telebot.types.BotCommand('/add_questions','📝 Testga savol qo\'shish'),
    telebot.types.BotCommand('/stop',         '🛑 Jarayonni to\'xtatish'),
    telebot.types.BotCommand('/stat',         '📊 Statistika'),
    telebot.types.BotCommand('/language',     '🌐 Tilni o\'zgartirish'),
])

database.init_db()

# Xotirada saqlash
user_sessions = {} # Test yaratish uchun
taking_quiz_sessions = {} # Yakkaxon test yechish uchun
group_sessions = {} # Guruhda test yechish uchun

def is_user_admin_or_starter(chat_id, user_id, starter_id):
    if user_id == starter_id:
        return True
    try:
        member = bot.get_chat_member(chat_id, user_id)
        if member.status in ['administrator', 'creator']:
            return True
    except Exception as e:
        print(f"Error checking chat member: {e}")
    return False

def is_group_quiz_active_and_intercept(message):
    chat_id = message.chat.id
    if message.chat.type in ['group', 'supergroup'] and chat_id in group_sessions:
        if message.text:
            first_word = message.text.split()[0].lower()
            command = first_word.split('@')[0]
            if command == '/stop':
                return False
        bot.reply_to(message, "⚠️ Guruhda hozirda faol test ketmoqda! Uni faqat guruh adminlari yoki boshlagan odam /stop orqali to'xtata oladi.")
        return True
    return False

def get_current_quiz_id(user_id):
    session = user_sessions.get(user_id)
    if session and 'quiz_id' in session:
        return session['quiz_id']
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT quiz_id FROM quizzes WHERE creator_id = ? ORDER BY quiz_id DESC LIMIT 1", (user_id,))
    res = cursor.fetchone()
    conn.close()
    if res:
        if user_id not in user_sessions:
            user_sessions[user_id] = {}
        user_sessions[user_id]['quiz_id'] = res[0]
        return res[0]
    return None

def get_user_language(user_id):
    try:
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        cursor.execute("SELECT language FROM users WHERE user_id = ?", (user_id,))
        res = cursor.fetchone()
        conn.close()
        return res[0] if (res and res[0]) else 'uz'
    except Exception:
        return 'uz'

def get_language_keyboard():
    markup = InlineKeyboardMarkup()
    markup.row(InlineKeyboardButton("O'zbekcha 🇺🇿", callback_data="set_lang_uz"),
               InlineKeyboardButton("English 🇬🇧", callback_data="set_lang_en"))
    markup.row(InlineKeyboardButton("Русский 🇷🇺", callback_data="set_lang_ru"),
               InlineKeyboardButton("Türkçe 🇹🇷", callback_data="set_lang_tr"))
    markup.row(InlineKeyboardButton("Deutsch 🇩🇪", callback_data="set_lang_de"),
               InlineKeyboardButton("العربية 🇸🇦", callback_data="set_lang_ar"))
    return markup

T = {
    'welcome_msg': {
        'uz': "Assalomu alaykum, {first_name}! 👋\n\nMen Telegram Quiz Botman. Men yordamida siz turli xil testlar va viktorinalar yaratishingiz mumkin.\n\nQuyidagi menyudan kerakli bo'limni tanlang:",
        'en': "Hello, {first_name}! 👋\n\nI am Telegram Quiz Bot. With me, you can create various quizzes and tests.\n\nPlease select the desired section from the menu below:",
        'ru': "Здравствуйте, {first_name}! 👋\n\nЯ Telegram Quiz Bot. С моей помощью вы можете создавать различные тесты и викторины.\n\nПожалуйста, выберите нужный раздел из меню ниже:",
        'tr': "Merhaba, {first_name}! 👋\n\nBen Telegram Quiz Bot. Benimle çeşitli testler va bilgi yarışmaları oluşturabilirsiniz.\n\nLütfen aşağıdaki menüden istediğiniz bölümü seçin:",
        'de': "Hallo, {first_name}! 👋\n\nIch bin der Telegram-Quiz-Bot. Mit mir können Sie verschiedene Quizze und Tests erstellen.\n\nBitte wählen Sie den gewünschten Bereich aus dem folgenden Menü:",
        'ar': "مرحباً، {first_name}! 👋\n\nأنا بوت الاختبارات. بمساعدتي، يمكنك إنشاء العديد من الاختبارات والمسابقات.\n\nالرجاء اختيار القسم المطلوب من القائمة أدناه:"
    },
    'group_not_supported_newquiz': {
        'uz': "Testni faqat shaxsiy chatda yaratish mumkin.",
        'en': "Quizzes can only be created in private chats.",
        'ru': "Тесты можно создавать только в личных чатах.",
        'tr': "Testler sadece özel sohbetlerde oluşturulabilir.",
        'de': "Quizze können nur in privaten Chats erstellt werden.",
        'ar': "يمكن إنشاء الاختبارات فقط في المحادثات الخاصة."
    },
    'waiting_for_title_prompt': {
        'uz': "Yangi test yaratamiz! 🎉\n\nIltimos, testning nomini yuboring (masalan: 'Matematika fanidan testlar').",
        'en': "Let's create a new quiz! 🎉\n\nPlease send the title of the quiz (e.g. 'Mathematics Test').",
        'ru': "Давайте создадим новый тест! 🎉\n\nПожалуйста, отправьте название теста (например: 'Тест по математике').",
        'tr': "Yeni bir test oluşturalım! 🎉\n\nLütfen testin başlığını gönderin (örneğin: 'Matematik Testi').",
        'de': "Lassen Sie uns ein neues Quiz erstellen! 🎉\n\nBitte senden Sie den Titel des Quiz (z. B. 'Mathematik-Test').",
        'ar': "دعنا ننشئ اختباراً جديداً! 🎉\n\nالرجاء إرسال عنوان الاختبار (مثال: 'اختبار الرياضيات')."
    },
    'waiting_for_description_prompt': {
        'uz': "Ajoyib! Endi test uchun qisqacha izoh (description) yuboring.\n\nAgar izoh qo'shishni xohlamasangiz, /skip buyrug'ini bosing.",
        'en': "Great! Now send a short description for the quiz.\n\nIf you don't want to add a description, click /skip.",
        'ru': "Отлично! Теперь отправьте краткое описание теста.\n\nЕсли вы не хотите добавлять описание, нажмите /skip.",
        'tr': "Harika! Şimdi test için kısa bir açıklama gönderin.\n\nAçıklama eklemek istemiyorsanız, /skip komutuna tıklayın.",
        'de': "Großartig! Senden Sie nun eine kurze Beschreibung für das Quiz.\n\nWenn Sie keine Beschreibung hinzufügen möchten, klicken Sie auf /skip.",
        'ar': "رائع! الآن أرسل وصفاً قصيراً للاختبار.\n\nإذا كنت لا تريد إضافة وصف، اضغط على /skip."
    },
    'description_saved': {
        'uz': "Izoh saqlandi.",
        'en': "Description saved.",
        'ru': "Описание сохранено.",
        'tr': "Açıklama kaydedildi.",
        'de': "Beschreibung gespeichert.",
        'ar': "تم حفظ الوصف."
    },
    'description_skipped': {
        'uz': "Izoh o'tkazib yuborildi.",
        'en': "Description skipped.",
        'ru': "Описание пропущено.",
        'tr': "Açıklama atlandı.",
        'de': "Beschreibung übersprungen.",
        'ar': "تم تخطي الوصف."
    },
    'ask_for_questions_prompt': {
        'uz': "\n\nEndi test uchun savollarni yuboring. Pastdagi **«Savol yaratish»** tugmasini bosing va o'z savolingizni tayyorlang.\n\nBarcha savollarni yuborib bo'lgach, **«Tugatish»** tugmasini bosing.",
        'en': "\n\nNow send the questions for the quiz. Click the **\"Create Question\"** button below and prepare your question.\n\nAfter sending all questions, click the **\"Finish\"** button.",
        'ru': "\n\nТеперь отправьте вопросы для теста. Нажмите кнопку **«Создать вопрос»** ниже и подготовьте свой вопрос.\n\nПосле отправки всех вопросов нажмите кнопку **«Завершить»**.",
        'tr': "\n\nŞimdi test için soruları gönderin. Aşağıdaki **\"Soru Oluştur\"** butonuna tıklayın ve sorunuzu hazırlayın.\n\nTüm soruları gönderdikten sonra **\"Bitir\"** butonuna tıklayın.",
        'de': "\n\nSenden Sie nun die Fragen für das Quiz. Klicken Sie unten auf die Schaltfläche **„Frage erstellen“** und bereiten Sie Ihre Frage vor.\n\nKlicken Sie nach dem Senden aller Fragen auf die Schaltfläche **„Beenden“**.",
        'ar': "\n\nالرجاء إرسال أسئلة الاختبار. اضغط على زر **«إنشاء سؤال»** أدناه وقم بإعداد سؤالك.\n\nبعد إرسال جميع الأسئلة، اضغط على زر **«إنهاء»**."
    },
    'poll_not_quiz': {
        'uz': "Iltimos, so'rovnomani **Quiz Mode** (Viktorina rejimi) da yarating!",
        'en': "Please create the poll in **Quiz Mode**!",
        'ru': "Пожалуйста, создайте опрос в **Режиме викторины**!",
        'tr': "Lütfen anketi **Bilgi Yarışması Modunda** oluşturun!",
        'de': "Bitte erstellen Sie die Umfrage im **Quiz-Modus**!",
        'ar': "الرجاء إنشاء الاستبيان في **وضع الاختبار**!"
    },
    'question_received_count': {
        'uz': "Yaxshi. \"{title}\" testingizda hozirda {count} ta savol bor. Agar oxirgi savolda xatoga yo'l qo'ygan bo'lsangiz, /undo buyrug'ini yuborish orqali orqaga qaytarishingiz mumkin.\n\nKeyingi savolni yaratish uchun yana **«Savol yaratish»** tugmasini bosing yoki tugatish uchun **«Tugatish»** tugmasini bosing.",
        'en': "Good. \"{title}\" now has {count} questions. If you made a mistake on the last question, you can roll it back by sending the /undo command.\n\nTo create the next question, click **\"Create Question\"** again, or click **\"Finish\"** to complete.",
        'ru': "Хорошо. В вашем тесте «{title}» сейчас {count} вопросов. Если вы допустили ошибку в последнем вопросе, вы можете отменить его, отправив команду /undo.\n\nЧтобы создать следующий вопрос, снова нажмите **«Создать вопрос»**, или нажмите **«Завершить»** для окончания.",
        'tr': "Güzel. \"{title}\" testinizde şu anda {count} soru var. Son soruda bir hata yaptıysanız, /undo komutunu göndererek geri alabilirsiniz.\n\nYeni soru oluşturmak için tekrar **\"Soru Oluştur\"** butonuna tıklayın veya tamamlamak için **\"Bitir\"** butonuna tıklayın.",
        'de': "Gut. Ihr Quiz „{title}“ hat jetzt {count} Fragen. Wenn Sie bei der letzten Frage einen Fehler gemacht haben, können Sie diesen durch Senden des Befehls /undo rückgängig machen.\n\nUm die nächste Frage zu erstellen, klicken Sie erneut auf **„Frage erstellen“** oder auf **„Beenden“**, um den Vorgang abzuschließen.",
        'ar': "حسنًا. اختبارك \"{title}\" يحتوي الآن على {count} أسئلة. إذا ارتكبت خطأً في السؤال الأخير، يمكنك التراجع عنه بإرسال الأمر /undo.\n\nلإنشاء السؤال التالي، اضغط على **«إنشاء سؤال»** مرة أخرى، أو اضغط на **«إنهاء»** للإكمال."
    },
    'undo_not_in_progress': {
        'uz': "Siz hozir savol qo'shish jarayonida emassiz.",
        'en': "You are not currently adding questions.",
        'ru': "Вы сейчас не находитесь в процессе добавления вопросов.",
        'tr': "Şu anda soru ekleme sürecinde değilsiniz.",
        'de': "Sie befinden sich derzeit nicht im Prozess des Hinzufügens von Fragen.",
        'ar': "أنت لا تقوم بإضافة أسئلة حالياً."
    },
    'undo_quiz_not_found': {
        'uz': "Yaratilayotgan test topilmadi.",
        'en': "Active quiz not found.",
        'ru': "Создаваемый тест не найден.",
        'tr': "Oluşturulmakta olan test bulunamadı.",
        'de': "Das zu erstellende Quiz wurde nicht gefunden.",
        'ar': "لم يتم العثور على الاختبار النشط."
    },
    'undo_no_questions': {
        'uz': "\"{title}\" testingizda hali hech qanday savol yo'q, shuning uchun /undo qilib bo'lmaydi.",
        'en': "There are no questions in your \"{title}\" quiz yet, so you cannot undo.",
        'ru': "В вашем тесте «{title}» еще нет вопросов, поэтому отмена невозможна.",
        'tr': "\"{title}\" testinizde henüz hiç soru yok, bu yüzden geri alma işlemi yapılamaz.",
        'de': "Es gibt noch keine Fragen in Ihrem „{title}“-Quiz, daher kann /undo nicht ausgeführt werden.",
        'ar': "لا توجد أسئلة في اختبارك \"{title}\" بعد، لذا لا يمكن التراجع."
    },
    'undo_success': {
        'uz': "🗑 Oxirgi savol muvaffaqiyatli o'chirildi!\nO'chirilgan savol: \"{q_text}\"\n\nHozirda \"{title}\" testingizda {count} ta savol qoldi.\nYangi savol yaratish uchun pastdagi **«Savol yaratish»** tugmasini bosing.",
        'en': "🗑 Last question successfully deleted!\nDeleted question: \"{q_text}\"\n\nNow \"{title}\" quiz has {count} questions left.\nTo add a new question, click **\"Create Question\"** below.",
        'ru': "🗑 Последний вопрос успешно удален!\nУдаленный вопрос: «{q_text}»\n\nТеперь в вашем тесте «{title}» осталось {count} вопросов.\nЧтобы добавить новый вопрос, нажмите **«Создать вопрос»** ниже.",
        'tr': "🗑 Son soru başarıyla silindi!\nSilinen soru: \"{q_text}\"\n\nŞu anda \"{title}\" testinizde {count} soru kaldı.\nYeni soru eklemek için aşağıdaki **\"Soru Oluştur\"** butonuna tıklayın.",
        'de': "🗑 Die letzte Frage wurde erfolgreich gelöscht!\nGelöschte Frage: „{q_text}“\n\nJetzt hat Ihr Quiz „{title}“ noch {count} Fragen übrig.\nUm eine neue Frage hinzuzufügen, klicken Sie unten auf **„Frage erstellen“**.",
        'ar': "🗑 تم حذف السؤال الأخير بنجاح!\nالسؤال المحذوف: \"{q_text}\"\n\nالآن اختبارك \"{title}\" يتبقى فيه {count} أسئلة.\nلإضافة سؤال جديد، اضغط على **«إنشاء سؤال»** أدناه."
    },
    'add_questions_prompt': {
        'uz': "OK! Oxirgi yaratgan \"{title}\" testingizga yana savollar qo'shishingiz mumkin (hozirda {count} ta savol bor). Savollarni (Poll) yuboravering.\n\nTugatganingizdan so'ng yana «Tugatish» tugmasini bosing.",
        'en': "OK! You can add more questions to your last created quiz \"{title}\" (currently has {count} questions). Keep sending Polls.\n\nAfter you finish, click **\"Finish\"** again.",
        'ru': "ОК! Вы можете добавить больше вопросов в свой последний созданный тест «{title}» (сейчас в нем {count} вопросов). Отправляйте опросы.\n\nПосле завершения снова нажмите кнопку **«Завершить»**.",
        'tr': "Tamam! Son oluşturduğunuz \"{title}\" testine daha fazla soru ekleyebilirsiniz (şu anda {count} soru var). Anketleri göndermeye devam edin.\n\nBitirdikten sonra tekrar **\"Bitir\"** butonuna tıklayın.",
        'de': "OK! Sie können Ihrem zuletzt erstellten Quiz „{title}“ weitere Fragen hinzufügen (aktuell hat es {count} Fragen). Senden Sie einfach weitere Umfragen.\n\nWenn Sie fertig sind, klicken Sie erneut auf **„Beenden“**.",
        'ar': "حسناً! يمكنك إضافة المزيد من الأسئلة إلى اختبارك الأخير \"{title}\" (يحتوي حالياً على {count} أسئلة). استمر في إرسال الاستبيانات.\n\nبعد الانتهاء، اضغط على **«إنهاء»** مرة أخرى."
    },
    'no_active_process': {
        'uz': "Hozircha hech qanday faol jarayon yo'q.",
        'en': "There are no active processes at the moment.",
        'ru': "На данный момент нет активных процессов.",
        'tr': "Şu anda aktif bir süreç yok.",
        'de': "Derzeit gibt es keine aktiven Prozesse.",
        'ar': "لا توجد عمليات نشطة في الوقت الحالي."
    },
    'stop_solo_quiz': {
        'uz': "🛑 Yakkaxon test to'xtatildi.",
        'en': "🛑 Solo quiz stopped.",
        'ru': "🛑 Одиночный тест остановлен.",
        'tr': "🛑 Bireysel test durduruldu.",
        'de': "🛑 Solo-Quiz gestoppt.",
        'ar': "🛑 تم إيقاف الاختبار الفردي."
    },
    'stop_group_quiz': {
        'uz': "🛑 Guruhdagi test to'xtatildi.",
        'en': "🛑 Group quiz stopped.",
        'ru': "🛑 Групповой тест остановлен.",
        'tr': "🛑 Grup testi durduruldu.",
        'de': "🛑 Gruppen-Quiz gestoppt.",
        'ar': "🛑 تم إيقاف اختبار المجموعة."
    },
    'stop_creation': {
        'uz': "🛑 Test yaratish jarayoni to'xtatildi.",
        'en': "🛑 Quiz creation process stopped.",
        'ru': "🛑 Процесс создания теста остановлен.",
        'tr': "🛑 Test oluşturma süreci durduruldu.",
        'de': "🛑 Quiz-Erstellungsprozess gestoppt.",
        'ar': "🛑 تم إيقاف عملية إنشاء الاختبار."
    },
    'stop_group_unauthorized': {
        'uz': "⚠️ Testni faqat adminlar yoki testni boshlagan odam to'xtata oladi!",
        'en': "⚠️ Only administrators or the user who started the quiz can stop it!",
        'ru': "⚠️ Только администраторы или пользователь, начавший тест, могут его остановить!",
        'tr': "⚠️ Sadece yöneticiler veya testi başlatan kullanıcı testi durdurabilir!",
        'de': "⚠️ Nur Administratoren oder der Benutzer, der das Quiz gestartet hat, können es stoppen!",
        'ar': "⚠️ يمكن فقط للمشرفين أو المستخدم الذي بدأ الاختبار إيقافه!"
    },
    'ask_for_time_limit': {
        'uz': "Iltimos, har bir savol uchun qancha vaqt berilishini tanlang:",
        'en': "Please choose the time limit for each question:",
        'ru': "Пожалуйста, выберите лимит времени для каждого вопроса:",
        'tr': "Lütfen her soru için zaman sınırını seçin:",
        'de': "Bitte wählen Sie das Zeitlimit für jede Frage:",
        'ar': "الرجاء اختيار الوقت المحدد لكل سؤال:"
    },
    'not_creating_quiz': {
        'uz': "Siz hozir hech qanday test yaratmayapsiz. Boshlash uchun /newquiz ni bosing.",
        'en': "You are not currently creating a quiz. Press /newquiz to start.",
        'ru': "Вы сейчас не создаете тест. Нажмите /newquiz, чтобы начать.",
        'tr': "Şu anda bir test oluşturmuyorsunuz. Başlamak için /newquiz komutuna basın.",
        'de': "Sie erstellen derzeit kein Quiz. Drücken Sie /newquiz, um zu starten.",
        'ar': "أنت لا تقوم بإنشاء اختبار حالياً. اضغط على /newquiz للبدء."
    },
    'time_limit_set': {
        'uz': "⏳ Vaqt {time_limit} soniya qilib belgilandi.\n\nSavollar va variantlar aralashtirilsinmi?",
        'en': "⏳ Time limit set to {time_limit} seconds.\n\nShould the questions and options be shuffled?",
        'ru': "⏳ Лимит времени установлен на {time_limit} сек.\n\nДолжны ли вопросы и варианты быть перемешаны?",
        'tr': "⏳ Zaman sınırı {time_limit} saniye olarak belirlendi.\n\nSorular ve seçenekler karıştırılsın mı?",
        'de': "⏳ Zeitlimit auf {time_limit} Sekunden festgelegt.\n\nSollen die Fragen und Optionen gemischt werden?",
        'ar': "⏳ تم تحديد الوقت بـ {time_limit} ثانية.\n\nهل تريد خلط الأسئلة والخيارات؟"
    },
    'time_limit_error': {
        'uz': "Xatolik yuz berdi. (Test topilmadi)",
        'en': "An error occurred. (Quiz not found)",
        'ru': "Произошла ошибка. (Тест не найден)",
        'tr': "Bir hata oluştu. (Test bulunamadı)",
        'de': "Ein Fehler ist aufgetreten. (Quiz nicht gefunden)",
        'ar': "حدث خطأ. (لم يتم العثور على الاختبار)"
    },
    'button_disabled': {
        'uz': "Bu tugma hozir ishlamaydi.",
        'en': "This button is not active at the moment.",
        'ru': "Эта кнопка сейчас не работает.",
        'tr': "Bu buton şu anda aktif değil.",
        'de': "Diese Schaltfläche ist derzeit nicht aktiv.",
        'ar': "هذا الزر لا يعمل حالياً."
    },
    'shuffle_saved': {
        'uz': "Aralashtirish rejimi saqlandi!",
        'en': "Shuffle mode saved!",
        'ru': "Режим перемешивания сохранен!",
        'tr': "Karıştırma modu kaydedildi!",
        'de': "Mischmodus gespeichert!",
        'ar': "تم حفظ وضع الخلط!"
    },
    'quiz_creation_success': {
        'uz': "✅ Test muvaffaqiyatli yaratildi va saqlandi!\n\nSizning ushbu \"{title}\" testingizda {count} ta test (savol) bo'ldi. 🎉\n\nQuyidagi tugmalar orqali testni o'zingiz yechishingiz yoki guruhlarga yuborishingiz mumkin:",
        'en': "✅ Quiz successfully created and saved!\n\nYour quiz \"{title}\" now has {count} questions. 🎉\n\nYou can use the buttons below to play solo or send it to groups:",
        'ru': "✅ Тест успешно создан и сохранен!\n\nВ вашем тесте «{title}» теперь {count} вопросов. 🎉\n\nВы можете использовать кнопки ниже, чтобы пройти его самостоятельно или отправить в группы:",
        'tr': "✅ Test başarıyla oluşturuldu ve kaydedildi!\n\n\"{title}\" testinizde şu anda {count} soru var. 🎉\n\nAşağıdaki butonları kullanarak testi kendiniz çözebilir veya gruplara gönderebilirsiniz:",
        'de': "✅ Quiz erfolgreich erstellt und gespeichert!\n\nIhr Quiz „{title}“ hat jetzt {count} Fragen. 🎉\n\nSie können die folgenden Schaltflächen verwenden, um es solo zu spielen oder an Gruppen zu senden:",
        'ar': "✅ تم إنشاء الاختبار وحفظه بنجاح!\n\nاختبارك \"{title}\" يحتوي الآن على {count} أسئلة. 🎉\n\nيمكنك استخدام الأزرار أدناه للحل الفردي أو إرساله للمجموعات:"
    },
    'play_solo': {
        'uz': "Yakkaxon yechish 👤",
        'en': "Play Solo 👤",
        'ru': "Решить соло 👤",
        'tr': "Bireysel Çöz 👤",
        'de': "Solo spielen 👤",
        'ar': "حل فردي 👤"
    },
    'play_group': {
        'uz': "Guruhda yechish 👥",
        'en': "Play in Group 👥",
        'ru': "Решить в группе 👥",
        'tr': "Grupta Çöz 👥",
        'de': "In Gruppe spielen 👥",
        'ar': "حل في مجموعة 👥"
    },
    'share_friends': {
        'uz': "Do'stlarga ulashish ↗️",
        'en': "Share with friends ↗️",
        'ru': "Поделиться с друзьями ↗️",
        'tr': "Arkadaşlarla Paylaş ↗️",
        'de': "Mit Freunden teilen ↗️",
        'ar': "مشاركة مع الأصدقاء ↗️"
    },
    'no_quizzes_yet': {
        'uz': "Sizda hali yaratilgan testlar yo'q. Yangi test yarating.",
        'en': "You don't have any quizzes created yet. Create a new quiz.",
        'ru': "У вас еще нет созданных тестов. Создайте новый тест.",
        'tr': "Henüz oluşturulmuş bir testiniz yok. Yeni bir test oluşturun.",
        'de': "Sie haben noch keine Quizze erstellt. Erstellen Sie ein neues Quiz.",
        'ar': "ليس لديك أي اختبارات منشأة بعد. أنشئ اختباراً جديداً."
    },
    'my_quizzes_title': {
        'uz': "📂 **Sizning testlaringiz:**\n\n",
        'en': "📂 **Your quizzes:**\n\n",
        'ru': "📂 **Ваши тесты:**\n\n",
        'tr': "📂 **Testleriniz:**\n\n",
        'de': "📂 **Ihre Quizze:**\n\n",
        'ar': "📂 **اختباراتك:**\n\n"
    },
    'results_section_coming_soon': {
        'uz': "Natijalar bo'limi tez kunda ishga tushadi! 🔜",
        'en': "Results section is coming soon! 🔜",
        'ru': "Раздел результатов скоро появится! 🔜",
        'tr': "Sonuçlar bölümü yakında hizmetinizde olacak! 🔜",
        'de': "Der Ergebnisbereich ist in Kürze verfügbar! 🔜",
        'ar': "قسم النتائج سيكون متاحاً قريباً! 🔜"
    },
    'ai_limit_reached': {
        'uz': "❌ Sizning bugungi limitingiz tugadi. Ertaga fayl yoki rasmlaringizni qaytadan yuborib ko'ring.",
        'en': "❌ Your daily limit has been reached. Please try uploading your files or images again tomorrow.",
        'ru': "❌ Ваш ежедневный лимит исчерпан. Пожалуйста, попробуйте отправить файлы или изображения завтра.",
        'tr': "❌ Günlük limitinize ulaşıldı. Lütfen dosyalarınızı veya resimlerinizi yarın tekrar göndermeyi deneyin.",
        'de': "❌ Ihr tägliches Limit wurde erreicht. Bitte versuchen Sie morgen erneut, Ihre Dateien oder Bilder hochzuladen.",
        'ar': "❌ لقد نفد حدك اليومي. الرجاء محاولة إرسال الملفات أو الصور مرة أخرى غداً."
    },
    'ai_not_active': {
        'uz': "⚠️ **AI Quiz funksiyasi faol emas!**\n\nUshbu funksiyadan foydalanish uchun bot egasi `.env` fayliga `GEMINI_API_KEY` kalitini qo'shishi kerak.",
        'en': "⚠️ **AI Quiz function is not active!**\n\nTo use this feature, the bot owner must add the `GEMINI_API_KEY` to the `.env` file.",
        'ru': "⚠️ **Функция AI Quiz не активна!**\n\nЧтобы использовать эту функцию, владелец бота должен добавить `GEMINI_API_KEY` в файл `.env`.",
        'tr': "⚠️ **AI Quiz özelliği aktif değil!**\n\nBu özelliği kullanmak için bot sahibinin `.env` dosyasına `GEMINI_API_KEY` eklemesi gerekir.",
        'de': "⚠️ **AI-Quiz-Funktion ist nicht aktiv!**\n\nUm diese Funktion zu nutzen, muss der Bot-Besitzer den `GEMINI_API_KEY` in die `.env`-Datei eintragen.",
        'ar': "⚠️ **ميزة اختبار الذكاء الاصطناعي غير نشطة!**\n\nلاستخدام هذه الميزة، يجب على مالك البوت إضافة `GEMINI_API_KEY` إلى ملف `.env`."
    },
    'ai_wrong_format': {
        'uz': "❌ Kechirasiz, faqat `.pdf`, `.docx`, `.pptx` va `.txt` formatidagi fayllarni qabul qila olaman.",
        'en': "❌ Sorry, I can only accept `.pdf`, `.docx`, `.pptx` and `.txt` files.",
        'ru': "❌ Извините, я принимаю только файлы форматов `.pdf`, `.docx`, `.pptx` и `.txt`.",
        'tr': "❌ Üzgünüm, yalnızca `.pdf`, `.docx`, `.pptx` ve `.txt` dosyalarını kabul edebilirim.",
        'de': "❌ Entschuldigung, ich kann nur `.pdf`-, `.docx`-, `.pptx`- und `.txt`-Dateien akzeptieren.",
        'ar': "❌ عذراً، يمكنني قبول ملفات `.pdf` و `.docx` و `.pptx` و `.txt` فقط."
    },
    'ai_downloading': {
        'uz': "📥 **Hujjat yuklab olinmoqda va o'qilmoqda...** Iltimos, kuting ⏳",
        'en': "📥 **Downloading and reading document...** Please wait ⏳",
        'ru': "📥 **Документ скачивается и читается...** Пожалуйста, подождите ⏳",
        'tr': "📥 **Belge indiriliyor ve okunuyor...** Lütfen bekleyin ⏳",
        'de': "📥 **Dokument wird heruntergeladen und gelesen...** Bitte warten ⏳",
        'ar': "📥 **يتم تنزيل وقراءة المستند...** الرجاء الانتظار ⏳"
    },
    'ai_empty_text': {
        'uz': "❌ Fayldan matn o'qib bo'lmadi yoki undagi matn juda qisqa (kamida 20 ta belgi bo'lishi kerak).",
        'en': "❌ Could not read text from the file or the text is too short (must be at least 20 characters).",
        'ru': "❌ Не удалось прочитать текст из файла или текст слишком короткий (минимум 20 символов).",
        'tr': "❌ Dosyadan metin okunamadı veya metin çok kısa (en az 20 karakter olmalıdır).",
        'de': "❌ Text konnte nicht aus der Datei gelesen werden oder der Text ist zu kurz (mindestens 20 Zeichen).",
        'ar': "❌ لم يتمكن من قراءة النص من الملف أو النص قصير جداً (يجب أن يكون 20 حرفاً على الأكثر)."
    },
    'ai_file_success': {
        'uz': "📄 **Fayl muvaffaqiyatli o'qildi!**\nHujjat: `{file_name}`\nMatn hajmi: {char_count} ta belgi.\n{limit_text}\nUshbu fayl asosida nechta test savoli yaratmoqchisiz? Tanlang 👇",
        'en': "📄 **File successfully read!**\nDocument: `{file_name}`\nText size: {char_count} characters.\n{limit_text}\nHow many questions do you want to create based on this file? Select below 👇",
        'ru': "📄 **Файл успешно прочитан!**\nДокумент: `{file_name}`\nРазмер текста: {char_count} символов.\n{limit_text}\nСколько вопросов вы хотите создать на основе этого файла? Выберите ниже 👇",
        'tr': "📄 **Dosya başarıyla okundu!**\nBelge: `{file_name}`\nMetin boyutu: {char_count} karakter.\n{limit_text}\nBu dosyaya dayanarak kaç soru oluşturmak istersiniz? Aşağıdan seçin 👇",
        'de': "📄 **Datei erfolgreich gelesen!**\nDokument: `{file_name}`\nTextgröße: {char_count} Zeichen.\n{limit_text}\nWie viele Fragen möchten Sie auf der Grundlage dieser Datei erstellen? Unten auswählen 👇",
        'ar': "📄 **تم قراءة الملف بنجاح!**\nالمستند: `{file_name}`\nحجم النص: {char_count} حرف.\n{limit_text}\nكم عدد الأسئلة التي تريد إنشاؤها بناءً على هذا الملف؟ اختر أدناه 👇"
    },
    'ai_photo_success': {
        'uz': "🖼 **Rasm muvaffaqiyatli qabul qilindi!**\n{limit_text}\nUshbu rasm/skrinshotdagi ma'lumotlar asosida nechta test savoli yaratmoqchisiz? Tanlang 👇",
        'en': "🖼 **Image successfully received!**\n{limit_text}\nHow many questions do you want to create based on this image/screenshot? Select below 👇",
        'ru': "🖼 **Изображение успешно получено!**\n{limit_text}\nСколько вопросов вы хотите создать на основе этого изображения/скриншота? Выберите ниже 👇",
        'tr': "🖼 **Resim başarıyla alındı!**\n{limit_text}\nBu resme/ekran görüntüsüne dayanarak kaç soru oluşturmak istersiniz? Aşağıdan seçin 👇",
        'de': "🖼 **Bild erfolgreich empfangen!**\n{limit_text}\nWie viele Fragen möchten Sie auf der Grundlage dieses Bildes/Screenshots erstellen? Unten auswählen 👇",
        'ar': "🖼 **تم استقبال الصورة بنجاح!**\n{limit_text}\nكم عدد الأسئلة التي تريد إنشاؤها بناءً على هذه الصورة/لقطة الشاشة؟ اختر أدناه 👇"
    },
    'ai_limit_text_user': {
        'uz': "\n⚠️ *Eslatma: Siz bir kunda maksimal 3 tagacha fayl yoki rasm yubora olasiz. (Bugungi qolgan limitingiz: {remaining} ta)*\n",
        'en': "\n⚠️ *Note: You can upload a maximum of 3 files or images per day. (Remaining limit today: {remaining})*\n",
        'ru': "\n⚠️ *Примечание: Вы можете отправлять максимум 3 файла или изображения в день. (Оставшийся лимит на сегодня: {remaining})*\n",
        'tr': "\n⚠️ *Not: Günde en fazla 3 dosya veya resim gönderebilirsiniz. (Bugünkü kalan limitiniz: {remaining})*\n",
        'de': "\n⚠️ *Hinweis: Sie können maximal 3 Dateien oder Bilder pro Tag hochladen. (Verbleibendes Limit heute: {remaining})*\n",
        'ar': "\n⚠️ *ملاحظة: يمكنك إرسال 3 ملفات أو صور كحد أقصى يومياً. (الحد المتبقي اليوم: {remaining})*\n"
    },
    'ai_limit_text_admin': {
        'uz': "\n⭐ *Siz uchun AI limiti cheksiz!*\n",
        'en': "\n⭐ *AI limit is unlimited for you!*\n",
        'ru': "\n⭐ *Лимит AI для вас неограничен!*\n",
        'tr': "\n⭐ *Sizin için AI limiti sınırsız!*\n",
        'de': "\n⭐ *Das AI-Limit ist für Sie unbegrenzt!*\n",
        'ar': "\n⭐ *حد الذكاء الاصطناعي غير محدود بالنسبة لك!*\n"
    },
    'ai_questions_count_selected': {
        'uz': "🎯 Savollar soni: **{count} ta** tanlandi.\n\nEndi, har bir savol uchun qancha vaqt berilishini tanlang 👇",
        'en': "🎯 Number of questions: **{count}** selected.\n\nNow, select the time limit for each question 👇",
        'ru': "🎯 Количество вопросов: **{count}** выбрано.\n\nТеперь выберите лимит времени для каждого вопроса 👇",
        'tr': "🎯 Soru sayısı: **{count}** seçildi.\n\nŞimdi, her soru için zaman sınırını seçin 👇",
        'de': "🎯 Anzahl der Fragen: **{count}** ausgewählt.\n\nWählen Sie nun das Zeitlimit für jede Frage 👇",
        'ar': "🎯 عدد الأسئلة: **{count}** تم اختيارها.\n\nالآن، اختر الوقت المحدد لكل سؤال 👇"
    },
    'ai_processing': {
        'uz': "🤖 **AI ma'lumotlarni tahlil qilmoqda va test yaratmoqda...**\nSavollar soni: {count} ta\nVaqt limiti: {time_limit} soniya\n\nIltimos, kuting. Bu jarayon 10-30 soniya vaqt olishi mumkin ⏳",
        'en': "🤖 **AI is analyzing data and generating test...**\nQuestions count: {count}\nTime limit: {time_limit} seconds\n\nPlease wait. This process may take 10-30 seconds ⏳",
        'ru': "🤖 **ИИ анализирует данные и создает тест...**\nКоличество вопросов: {count}\nЛимит времени: {time_limit} секунд\n\nПожалуйста, подождите. Этот процесс может занять 10-30 секунд ⏳",
        'tr': "🤖 **AI verileri analiz ediyor ve test oluşturuyor...**\nSoru sayısı: {count}\nZaman sınırı: {time_limit} saniye\n\nLütfen bekleyin. Bu işlem 10-30 saniye sürebilir ⏳",
        'de': "🤖 **KI analysiert Daten und generiert Tests...**\nFragenanzahl: {count}\nZeitlimit: {time_limit} Sekunden\n\nBitte warten. Dieser Vorgang kann 10-30 Sekunden dauern ⏳",
        'ar': "🤖 **الذكاء الاصطناعي يحلل البيانات وينشئ الاختبار...**\nعدد الأسئلة: {count}\nالوقت المحدد: {time_limit} ثانية\n\nالرجاء الانتظار. قد تستغرق هذه العملية 10-30 ثانية ⏳"
    },
    'ai_creation_success': {
        'uz': "🎉 **AI Quiz muvaffaqiyatli yaratildi!**\n\n📌 **Sarlavha:** {title}\n📝 **Izoh:** {description}\n❓ **Savollar soni:** {count} ta\n\nQuyidagi tugmalar orqali testni yechishingiz yoki ulashishingiz mumkin:",
        'en': "🎉 **AI Quiz successfully created!**\n\n📌 **Title:** {title}\n📝 **Description:** {description}\n❓ **Questions count:** {count}\n\nYou can use the buttons below to play or share the quiz:",
        'ru': "🎉 **ИИ-викторина успешно создана!**\n\n📌 **Название:** {title}\n📝 **Описание:** {description}\n❓ **Количество вопросов:** {count}\n\nВы можете использовать кнопки ниже, чтобы пройти или поделиться тестом:",
        'tr': "🎉 **AI Testi başarıyla oluşturuldu!**\n\n📌 **Başlık:** {title}\n📝 **Açıklama:** {description}\n❓ **Soru sayısı:** {count}\n\nTesti çözmek veya paylaşmak için aşağıdaki butonları kullanabilirsiniz:",
        'de': "🎉 **AI-Quiz erfolgreich erstellt!**\n\n📌 **Titel:** {title}\n📝 **Beschreibung:** {description}\n❓ **Fragenanzahl:** {count}\n\nSie können die folgenden Schaltflächen verwenden, um das Quiz zu spielen oder zu teilen:",
        'ar': "🎉 **تم إنشاء اختبار الذكاء الاصطناعي بنجاح!**\n\n📌 **العنوان:** {title}\n📝 **الوصف:** {description}\n❓ **عدد الأسئلة:** {count}\n\nيمكنك استخدام الأزرار أدناه للحل أو مشاركة الاختبار:"
    },
    'ai_empty_response': {
        'uz': "❌ Xatolik: AI matndan ma'lumot ololmadi yoki javob qaytarmadi.",
        'en': "❌ Error: AI could not extract information from the text or returned an empty response.",
        'ru': "❌ Ошибка: ИИ не смог извлечь информацию из текста или вернул пустой ответ.",
        'tr': "❌ Hata: AI metinden bilgi çıkaramadı veya boş bir yanıt döndürdü.",
        'de': "❌ Fehler: KI konnte keine Informationen aus dem Text extrahieren oder hat eine leere Antwort zurückgegeben.",
        'ar': "❌ خطأ: لم يتمكن الذكاء الاصطناعي من استخراج المعلومات من النص أو أعاد استجابة فارغة."
    },
    'ai_no_questions': {
        'uz': "❌ Xatolik: AI birorta ham test savoli yarata olmadi. Boshqa fayl yoki rasm yuborib ko'ring.",
        'en': "❌ Error: AI could not generate any quiz questions. Please try sending another file or image.",
        'ru': "❌ Ошибка: ИИ не смог создать ни одного тестового вопроса. Попробуйте отправить другой файл или изображение.",
        'tr': "❌ Hata: AI hiç soru oluşturamadı. Lütfen başka bir dosya veya resim göndermeyi deneyin.",
        'de': "❌ Fehler: KI konnte keine Quizfragen generieren. Bitte versuchen Sie, eine andere Datei oder ein anderes Bild zu senden.",
        'ar': "❌ خطأ: لم يتمكن الذكاء الاصطناعي من إنشاء أي أسئلة. الرجاء محاولة إرسال ملف أو صورة أخرى."
    },
    'group_quiz_active_intercept_warning': {
        'uz': "⚠️ Guruhda hozirda faol test ketmoqda! Uni faqat guruh adminlari yoki boshlagan odam /stop orqali to'xtata oladi.",
        'en': "⚠️ An active quiz is currently running in this group! Only group admins or the user who started it can stop it with /stop.",
        'ru': "⚠️ В этой группе сейчас идет активный тест! Только администраторы группы или пользователь, начавший его, могут остановить его с помощью /stop.",
        'tr': "⚠️ Bu grupta şu anda aktif bir test var! Yalnızca grup yöneticileri veya testi başlatan kullanıcı /stop ile durdurabilir.",
        'de': "⚠️ In dieser Gruppe läuft derzeit ein aktives Quiz! Nur Gruppen-Admins oder der Benutzer, der es gestartet hat, können es mit /stop stoppen.",
        'ar': "⚠️ هناك اختبار نشط يجري حالياً في هذه المجموعة! يمكن فقط لمشرفي المجموعة أو المستخدم الذي بدأه إيقافه باستخدام /stop."
    },
    'solo_quiz_starting': {
        'uz': "🏁 \"{title}\" testi boshlanmoqda...\n\n❓ Savollar soni: {question_count} ta\n⏳ Har bir savol uchun belgilangan vaqt: {time_limit} soniya\n\nBoshladik! 🚀",
        'en': "🏁 Starting quiz \"{title}\"...\n\n❓ Number of questions: {question_count}\n⏳ Time limit per question: {time_limit} seconds\n\nLet's start! 🚀",
        'ru': "🏁 Начинается тест «{title}»...\n\n❓ Количество вопросов: {question_count}\n⏳ Лимит времени на вопрос: {time_limit} секунд\n\nПоехали! 🚀",
        'tr': "🏁 \"{title}\" testi başlıyor...\n\n❓ Soru sayısı: {question_count}\n⏳ Soru başına zaman sınırı: {time_limit} saniye\n\nBaşlayalım! 🚀",
        'de': "🏁 Starten des Quiz „{title}“...\n\n❓ Anzahl der Fragen: {question_count}\n⏳ Zeitlimit pro Frage: {time_limit} Sekunden\n\nLasst uns beginnen! 🚀",
        'ar': "🏁 يبدأ الاختبار \"{title}\"...\n\n❓ عدد الأسئلة: {question_count}\n⏳ الوقت المحدد لكل سؤال: {time_limit} ثانية\n\nلنبدأ! 🚀"
    },
    'quiz_leaderboard_finished': {
        'uz': "🏁 **Test yakunlandi!**\n\n📊 **Test nomi:** \"{title}\"\n👥 **Jami qatnashchilar:** {user_count} ta\n\n🏆 **NATIJALAR:**\n{list_text}",
        'en': "🏁 **Quiz completed!**\n\n📊 **Quiz name:** \"{title}\"\n👥 **Total participants:** {user_count}\n\n🏆 **RESULTS:**\n{list_text}",
        'ru': "🏁 **Тест завершен!**\n\n📊 **Название теста:** «{title}»\n👥 **Всего участников:** {user_count}\n\n🏆 **РЕЗУЛЬТАТЫ:**\n{list_text}",
        'tr': "🏁 **Test tamamlandı!**\n\n📊 **Test adı:** \"{title}\"\n👥 **Toplam katılımcı:** {user_count}\n\n🏆 **SONUÇLAR:**\n{list_text}",
        'de': "🏁 **Quiz abgeschlossen!**\n\n📊 **Quizname:** „{title}“\n👥 **Teilnehmer insgesamt:** {user_count}\n\n🏆 **ERGEBNISSE:**\n{list_text}",
        'ar': "🏁 **اكتمل الاختبار!**\n\n📊 **اسم الاختبار:** \"{title}\"\n👥 **إجمالي المشاركين:** {user_count}\n\n🏆 **النتائج:**\n{list_text}"
    },
    'quiz_finished_solo': {
        'uz': "🏁 **Test yakunlandi!**\n\n📊 Sizning natijangiz: **{score} / {total}** ta to'g'ri javob. 🎉",
        'en': "🏁 **Quiz completed!**\n\n📊 Your score: **{score} / {total}** correct answers. 🎉",
        'ru': "🏁 **Тест завершен!**\n\n📊 Ваш результат: **{score} / {total}** правильных ответов. 🎉",
        'tr': "🏁 **Test tamamlandı!**\n\n📊 Skorunuz: **{score} / {total}** doğru cevap. 🎉",
        'de': "🏁 **Quiz abgeschlossen!**\n\n📊 Ihr Ergebnis: **{score} / {total}** richtige Antworten. 🎉",
        'ar': "🏁 **اكتمل الاختبار!**\n\n📊 نتيجتك: **{score} / {total}** إجابات صحيحة. 🎉"
    },
    'choose_language': {
        'uz': "🌐 Iltimos, o'zingizga qulay tilni tanlang:",
        'en': "🌐 Please select a language that is convenient for you:",
        'ru': "🌐 Пожалуйста, выберите удобный для вас язык:",
        'tr': "🌐 Lütfen sizin için uygun olan dili seçin:",
        'de': "🌐 Bitte wählen Sie eine für Sie passende Sprache:",
        'ar': "🌐 الرجاء اختيار لغة مناسبة لك:"
    },
    'solve_quiz': {
        'uz': "👉 Yechish",
        'en': "👉 Solve",
        'ru': "👉 Решить",
        'tr': "👉 Çöz",
        'de': "👉 Lösen",
        'ar': "👉 حل"
    },
    'main_menu': {
        'uz': "🏠 Asosiy menyu:",
        'en': "🏠 Main menu:",
        'ru': "🏠 Главное меню:",
        'tr': "🏠 Ana menü:",
        'de': "🏠 Hauptmenü:",
        'ar': "🏠 القائمة الرئيسية:"
    },
    'no_active_quiz': {
        'uz': "Sizda hali yaratilgan test yo'q.",
        'en': "You don't have any quizzes created yet.",
        'ru': "У вас еще нет созданных тестов.",
        'tr': "Henüz oluşturulmuş bir testiniz yok.",
        'de': "Sie haben noch keine Quizze erstellt.",
        'ar': "ليس لديك أي اختبارات منشأة بعد."
    }
}

def get_question_keyboard(lang='uz'):
    markup = ReplyKeyboardMarkup(resize_keyboard=True, row_width=1)
    labels = {
        'uz': ('Savol yaratish', 'Tugatish'),
        'en': ('Create Question', 'Finish'),
        'ru': ('Создать вопрос', 'Завершить'),
        'tr': ('Soru Oluştur', 'Bitir'),
        'de': ('Frage erstellen', 'Beenden'),
        'ar': ('إنشاء سؤال', 'إنهاء')
    }
    lbl_poll, lbl_done = labels.get(lang, labels['uz'])
    btn_poll = KeyboardButton(lbl_poll, request_poll=KeyboardButtonPollType(type='quiz'))
    btn_done = KeyboardButton(lbl_done)
    markup.add(btn_poll, btn_done)
    return markup

def get_main_keyboard(lang='uz'):
    markup = ReplyKeyboardMarkup(resize_keyboard=True, row_width=2)
    labels = {
        'uz': ('➕ Yangi test yaratish', '📂 Mening testlarim', '📊 Natijalar', '🌐 Tilni o\'zgartirish'),
        'en': ('➕ Create New Quiz', '📂 My Quizzes', '📊 Results', '🌐 Change Language'),
        'ru': ('➕ Создать новый тест', '📂 Мои тесты', '📊 Результаты', '🌐 Сменить язык'),
        'tr': ('➕ Yeni Test Oluştur', '📂 Testlerim', '📊 Sonuçlar', '🌐 Dili Değiştir'),
        'de': ('➕ Neuen Test erstellen', '📂 Meine Tests', '📊 Ergebnisse', '🌐 Sprache ändern'),
        'ar': ('➕ إنشاء اختبار جديد', '📂 اختباراتي', '📊 النتائج', '🌐 تغيير اللغة')
    }
    lbl_new, lbl_my, lbl_results, lbl_lang = labels.get(lang, labels['uz'])
    btn_new = KeyboardButton(lbl_new)
    btn_my = KeyboardButton(lbl_my)
    btn_results = KeyboardButton(lbl_results)
    btn_lang = KeyboardButton(lbl_lang)
    markup.add(btn_new, btn_my, btn_results, btn_lang)
    return markup

@bot.message_handler(commands=['start'])
def send_welcome(message):
    if is_group_quiz_active_and_intercept(message):
        return
    chat_type = message.chat.type
    text = message.text.split()
    payload = text[1] if len(text) > 1 else None

    # Agar guruhda bo'lsa
    if chat_type in ['group', 'supergroup']:
        if payload and payload.startswith('quiz_'):
            try:
                quiz_id = int(payload.split('_')[1])
                init_group_quiz(message.chat.id, quiz_id, message.from_user.id)
            except (IndexError, ValueError):
                bot.send_message(message.chat.id, "Xato link formatı.")
        return

    # Shaxsiy chat (Private)
    user_id = message.from_user.id
    username = message.from_user.username
    first_name = message.from_user.first_name
    
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    # Check if user exists
    cursor.execute("SELECT language FROM users WHERE user_id = ?", (user_id,))
    user_res = cursor.fetchone()
    
    if not user_res:
        # Register user with None language first
        cursor.execute("INSERT OR IGNORE INTO users (user_id, username, first_name, language) VALUES (?, ?, ?, NULL)", 
                       (user_id, username, first_name))
        conn.commit()
        user_lang = None
    else:
        user_lang = user_res[0]
    conn.close()

    # Deep linking (Yakkaxon test yechish)
    if payload and payload.startswith('quiz_'):
        # For deep linking, if user has no language, default to 'uz' to avoid blocking
        if user_lang is None:
            conn = sqlite3.connect(database.DB_PATH)
            cursor = conn.cursor()
            cursor.execute("UPDATE users SET language = 'uz' WHERE user_id = ?", (user_id,))
            conn.commit()
            conn.close()
            user_lang = 'uz'
        try:
            quiz_id = int(payload.split('_')[1])
            start_taking_quiz(message, quiz_id)
        except (IndexError, ValueError):
            bot.reply_to(message, "Xato link formati.")
        return

    database.set_user_state(user_id, 'none')
    
    if user_lang is None:
        # User needs to select a language first!
        welcome_prompt = (
            "Iltimos, botdan foydalanish uchun tilni tanlang:\n"
            "Please select a language to use the bot:\n"
            "Пожалуйста, выберите язык для использования бота:\n"
            "Lütfen botu kullanmak için bir dil seçin:\n"
            "Bitte wählen Sie eine Sprache, um den Bot zu verwenden:\n"
            "الرجاء اختيار لغة لاستخدام البوت:"
        )
        bot.reply_to(message, welcome_prompt, reply_markup=get_language_keyboard())
    else:
        # Show welcome in selected language
        msg_template = T['welcome_msg'].get(user_lang, T['welcome_msg']['uz'])
        welcome_text = msg_template.format(first_name=first_name)
        bot.reply_to(message, welcome_text, reply_markup=get_main_keyboard(user_lang))

@bot.callback_query_handler(func=lambda call: call.data.startswith("set_lang_"))
def handle_set_language(call):
    user_id = call.from_user.id
    lang = call.data.split('_')[2]
    
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    cursor.execute("UPDATE users SET language = ? WHERE user_id = ?", (lang, user_id))
    conn.commit()
    conn.close()
    
    database.set_user_state(user_id, 'none')
    
    confirms = {
        'uz': "✅ Til tanlandi: O'zbekcha 🇺🇿\n\nAssalomu alaykum! Bosh menyuga xush kelibsiz.",
        'en': "✅ Language selected: English 🇬🇧\n\nHello! Welcome to the main menu.",
        'ru': "✅ Язык выбран: Русский 🇷🇺\n\nЗдравствуйте! Добро пожаловать в главное меню.",
        'tr': "✅ Dil seçildi: Türkçe 🇹🇷\n\nMerhaba! Ana menüye hoş geldiniz.",
        'de': "✅ Sprache ausgewählt: Deutsch 🇩🇪\n\nHallo! Willkommen im Hauptmenü.",
        'ar': "✅ تم اختيار اللغة: العربية 🇸🇦\n\nمرحباً! أهلاً بك في القائمة الرئيسية."
    }
    
    bot.answer_callback_query(call.id, confirms.get(lang, confirms['uz']))
    
    try:
        bot.delete_message(call.message.chat.id, call.message.message_id)
    except Exception:
        pass
    
    welcome_texts = {
        'uz': "Quyidagi menyudan kerakli bo'limni tanlang:",
        'en': "Please select the desired section from the menu below:",
        'ru': "Пожалуйста, выберите нужный раздел из меню ниже:",
        'tr': "Lütfen aşağıdaki menüden istediğiniz bölümü seçin:",
        'de': "Bitte wählen Sie den gewünschten Bereich aus dem folgenden Menü:",
        'ar': "الرجاء اختيار القسم المطلوب من القائمة أدناه:"
    }
    
    bot.send_message(call.message.chat.id, welcome_texts.get(lang, welcome_texts['uz']), reply_markup=get_main_keyboard(lang))

def prepare_questions(quiz_id):
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT shuffle_mode FROM quizzes WHERE quiz_id = ?", (quiz_id,))
    res = cursor.fetchone()
    shuffle_mode = res[0] if res else 0

    cursor.execute("SELECT question_id, question_text, options, correct_option_id, explanation, time_limit FROM questions WHERE quiz_id = ?", (quiz_id,))
    questions = cursor.fetchall()
    conn.close()
    
    questions = [list(q) for q in questions]

    if shuffle_mode in [1, 3]: # Savollarni aralashtirish
        random.shuffle(questions)
        
    for i in range(len(questions)):
        q_id, q_text, options_json, correct_id, exp, time_limit = questions[i]
        options = json.loads(options_json)
        
        if shuffle_mode in [2, 3]: # Variantlarni aralashtirish
            correct_text = options[correct_id]
            random.shuffle(options)
            new_correct_id = options.index(correct_text)
            
            questions[i][2] = json.dumps(options)
            questions[i][3] = new_correct_id
            
    return questions

# --- YAKKAXON TEST YECHISH ---
def start_taking_quiz(message, quiz_id):
    user_id = message.from_user.id
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT title, description FROM quizzes WHERE quiz_id = ?", (quiz_id,))
    quiz = cursor.fetchone()
    conn.close()
    if not quiz:
        all_q = database.get_all_quizzes()
        all_ids = [str(q[0]) for q in all_q]
        bot.reply_to(message, f"Kechirasiz, bunday test (ID: {quiz_id}) topilmadi.\nBazadagi mavjud testlar IDlari: {', '.join(all_ids)}")
        return
        
    questions = prepare_questions(quiz_id)
    
    if not questions:
        bot.reply_to(message, "Bu testda hali savollar yo'q.")
        return
        
    title, desc = quiz
    desc_text = f"\nIzoh: {desc}" if desc else ""
    bot.reply_to(message, f"🎯 **{title}** testini boshlaymiz!{desc_text}\n\nJami savollar: {len(questions)}\nOmad!", parse_mode='Markdown')
    
    taking_quiz_sessions[user_id] = {
        'quiz_id': quiz_id,
        'questions': questions,
        'current_q_index': 0,
        'score': 0
    }
    send_next_question(user_id, message.chat.id)

def send_next_question(user_id, chat_id):
    session = taking_quiz_sessions.get(user_id)
    if not session: return
        
    idx = session['current_q_index']
    questions = session['questions']
    
    if idx < len(questions):
        q = questions[idx]
        q_id, q_text, options_json, correct_id, exp, time_limit = q
        options = json.loads(options_json)
        
        t_limit = time_limit if time_limit else 15

        msg = bot.send_poll(
            chat_id, 
            question=f"[{idx+1}/{len(questions)}] {q_text}", 
            options=options, 
            type='quiz', 
            correct_option_id=correct_id,
            explanation=exp,
            is_anonymous=False,
            open_period=t_limit
        )
        session['current_poll_id'] = msg.poll.id
        session['correct_option_id'] = correct_id
        session['chat_id'] = chat_id
        
        # Vaqt tugaguncha javob bermasa avtomatik o'tish uchun taymer
        if 'timer' in session and session['timer']:
            session['timer'].cancel()
            
        session['timer'] = threading.Timer(t_limit + 1.0, handle_question_timeout, args=[user_id, idx])
        session['timer'].start()
    else:
        score = session['score']
        total = len(questions)
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        cursor.execute("INSERT INTO results (user_id, quiz_id, score) VALUES (?, ?, ?)", (user_id, session['quiz_id'], score))
        conn.commit()
        conn.close()
        
        markup = InlineKeyboardMarkup()
        start_url = f"https://t.me/{bot_info.username}?start=quiz_{session['quiz_id']}"
        group_url = f"https://t.me/{bot_info.username}?startgroup=quiz_{session['quiz_id']}"
        markup.add(InlineKeyboardButton("Yakkaxon yechish 👤", url=start_url))
        markup.add(InlineKeyboardButton("Guruhda yechish 👥", url=group_url))
        markup.add(InlineKeyboardButton("Do'stlarga ulashish ↗️", url=f"https://t.me/share/url?url={start_url}"))

        bot.send_message(chat_id, f"🏁 Test yakunlandi!\n\nSizning natijangiz: {score} / {total}\nBarakalla! 🎉", reply_markup=markup)
        del taking_quiz_sessions[user_id]

def handle_question_timeout(user_id, q_index):
    session = taking_quiz_sessions.get(user_id)
    if session and session['current_q_index'] == q_index:
        session['current_q_index'] += 1
        send_next_question(user_id, session['chat_id'])

def init_group_quiz(chat_id, quiz_id, starter_id):
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT title FROM quizzes WHERE quiz_id = ?", (quiz_id,))
    quiz = cursor.fetchone()
    conn.close()
    
    if not quiz:
        all_q = database.get_all_quizzes()
        all_ids = [str(q[0]) for q in all_q]
        bot.send_message(chat_id, f"Kechirasiz, bunday test (ID: {quiz_id}) topilmadi.\nBazadagi mavjud testlar IDlari: {', '.join(all_ids)}")
        return
        
    questions = prepare_questions(quiz_id)
    
    if not questions:
        bot.send_message(chat_id, "Bu testda hali savollar yo'q.")
        return
        
    group_sessions[chat_id] = {
        'quiz_id': quiz_id,
        'title': quiz[0],
        'questions': questions,
        'current_q_index': 0,
        'participants': {},
        'status': 'waiting',
        'current_poll_id': None,
        'correct_option_id': None,
        'starter_id': starter_id
    }
    
    markup = InlineKeyboardMarkup()
    markup.add(InlineKeyboardButton("Men tayyorman! (0)", callback_data="ready"))
    markup.add(InlineKeyboardButton("▶️ Boshlash", callback_data="start_group_quiz"))
    
    msg = bot.send_message(chat_id, f"📢 **{quiz[0]}** testi boshlanmoqda!\n\nQatnashish uchun pastdagi tugmani bosing.", reply_markup=markup, parse_mode='Markdown')
    group_sessions[chat_id]['message_id'] = msg.message_id

@bot.callback_query_handler(func=lambda call: call.data in ["ready", "start_group_quiz"])
def handle_group_callbacks(call):
    chat_id = call.message.chat.id
    session = group_sessions.get(chat_id)
    if not session or session['status'] != 'waiting':
        bot.answer_callback_query(call.id, "Test hozir bu holatda emas.")
        return
        
    if call.data == "ready":
        user_id = call.from_user.id
        first_name = call.from_user.first_name
        
        if user_id not in session['participants']:
            session['participants'][user_id] = {'first_name': first_name, 'score': 0}
            count = len(session['participants'])
            
            markup = InlineKeyboardMarkup()
            markup.add(InlineKeyboardButton(f"Men tayyorman! ({count})", callback_data="ready"))
            markup.add(InlineKeyboardButton("▶️ Boshlash", callback_data="start_group_quiz"))
            
            bot.edit_message_reply_markup(chat_id, call.message.message_id, reply_markup=markup)
            bot.answer_callback_query(call.id, "Siz ro'yxatga olindingiz!")
        else:
            bot.answer_callback_query(call.id, "Siz allaqachon ro'yxatdan o'tgansiz.")
            
    elif call.data == "start_group_quiz":
        if len(session['participants']) == 0:
            session['participants'][call.from_user.id] = {'first_name': call.from_user.first_name, 'score': 0}
            
        session['status'] = 'playing'
        bot.edit_message_text(f"🚀 **{session['title']}** testi boshlandi!\nDiqqat qiling, savollar kelmoqda...", chat_id=chat_id, message_id=session['message_id'], parse_mode='Markdown')
        bot.answer_callback_query(call.id, "Test boshlandi!")
        
        timer = threading.Timer(3.0, send_next_group_question, args=[chat_id])
        session['timer'] = timer
        timer.start()

def send_next_group_question(chat_id):
    session = group_sessions.get(chat_id)
    if not session or session['status'] != 'playing': return
    
    idx = session['current_q_index']
    questions = session['questions']
    
    if idx < len(questions):
        q = questions[idx]
        q_id, q_text, options_json, correct_id, exp, time_limit = q
        options = json.loads(options_json)
        
        t_limit = time_limit if time_limit else 15

        msg = bot.send_poll(
            chat_id, 
            question=f"[{idx+1}/{len(questions)}] {q_text}", 
            options=options, 
            type='quiz', 
            correct_option_id=correct_id,
            explanation=exp,
            is_anonymous=False,
            open_period=t_limit
        )
        session['current_poll_id'] = msg.poll.id
        session['correct_option_id'] = correct_id
        session['current_q_index'] += 1
        
        timer = threading.Timer(t_limit + 1.0, send_next_group_question, args=[chat_id])
        session['timer'] = timer
        timer.start()
    else:
        session['status'] = 'finished'
        if 'timer' in session and session['timer']:
            session['timer'].cancel()
        participants = session['participants'].values()
        sorted_p = sorted(participants, key=lambda x: x['score'], reverse=True)
        
        text = f"🏁 **{session['title']}** testi yakunlandi!\n\n🏆 **NATIJALAR (Leaderboard):**\n\n"
        for i, p in enumerate(sorted_p):
            medal = "🥇" if i == 0 else "🥈" if i == 1 else "🥉" if i == 2 else "🔹"
            text += f"{medal} {p['first_name']} - {p['score']} ta to'g'ri\n"
            
        markup = InlineKeyboardMarkup()
        start_url = f"https://t.me/{bot_info.username}?start=quiz_{session['quiz_id']}"
        group_url = f"https://t.me/{bot_info.username}?startgroup=quiz_{session['quiz_id']}"
        markup.add(InlineKeyboardButton("Yakkaxon yechish 👤", url=start_url))
        markup.add(InlineKeyboardButton("Guruhda yechish 👥", url=group_url))
        markup.add(InlineKeyboardButton("Do'stlarga ulashish ↗️", url=f"https://t.me/share/url?url={start_url}"))

        bot.send_message(chat_id, text, parse_mode='Markdown', reply_markup=markup)
        del group_sessions[chat_id]


# --- POLL JAVOBLARINI USHLASH (Yakkaxon va Guruh) ---
@bot.poll_answer_handler()
def handle_poll_answer(poll_answer):
    user_id = poll_answer.user.id
    poll_id = poll_answer.poll_id
    selected_option = poll_answer.option_ids[0] if poll_answer.option_ids else None
    
    # 1. Yakkaxon testni tekshirish
    session = taking_quiz_sessions.get(user_id)
    if session and session.get('current_poll_id') == poll_id:
        if 'timer' in session and session['timer']:
            session['timer'].cancel()
            
        if selected_option == session['correct_option_id']:
            session['score'] += 1
            
        session['current_q_index'] += 1
        send_next_question(user_id, session['chat_id'])
        return
        
    # 2. Guruh testini tekshirish
    for chat_id, g_session in group_sessions.items():
        if g_session.get('current_poll_id') == poll_id:
            if user_id not in g_session['participants']:
                g_session['participants'][user_id] = {'first_name': poll_answer.user.first_name, 'score': 0}

            if selected_option == g_session['correct_option_id']:
                g_session['participants'][user_id]['score'] += 1
            break


# --- TEST YARATISH ---
@bot.message_handler(commands=['newquiz'])
def create_new_quiz(message):
    if is_group_quiz_active_and_intercept(message):
        return
    user_id = message.from_user.id
    lang = get_user_language(user_id)
    if message.chat.type in ['group', 'supergroup']:
        bot.reply_to(message, T['group_not_supported_newquiz'].get(lang, T['group_not_supported_newquiz']['uz']))
        return
    database.set_user_state(user_id, 'waiting_for_title')
    bot.reply_to(message, T['waiting_for_title_prompt'].get(lang, T['waiting_for_title_prompt']['uz']), reply_markup=ReplyKeyboardRemove())

@bot.message_handler(commands=['fix_my_id'])
def fix_id(message):
    if is_group_quiz_active_and_intercept(message):
        return
    user_id = message.from_user.id
    # Bazadagi Muhammadyusuf ID-si: 6559589296
    old_id = 6559589296
    database.fix_user_quizzes(old_id, user_id)
    bot.reply_to(message, f"✅ Barcha eski testlaringiz yangi ID-ingizga ({user_id}) biriktirildi! Endi 'Mening testlarim' bo'limini tekshirib ko'ring.")

@bot.message_handler(commands=['debug_me'])
def debug_me(message):
    if is_group_quiz_active_and_intercept(message):
        return
    user_id = message.from_user.id
    my_quizzes = database.get_my_quizzes(user_id)
    all_quizzes = database.get_all_quizzes()
    total_users = database.get_users_count()
    total_results = database.get_results_count()
    text = (
        f"🔍 **DEBUG MA'LUMOTLARI:**\n"
        f"Sizning ID: `{user_id}`\n"
        f"Sizning testlaringiz soni: {len(my_quizzes)}\n"
        f"Bazadagi jami testlar soni: {len(all_quizzes)}\n"
        f"Jami foydalanuvchilar soni: {total_users} ta\n"
        f"Jami yechilgan testlar soni: {total_results} ta\n"
        f"Baza fayli: `{database.DB_PATH}`\n"
    )
    bot.reply_to(message, text, parse_mode='Markdown')

@bot.message_handler(commands=['stat', 'stats'])
def show_statistics(message):
    if is_group_quiz_active_and_intercept(message):
        return
    try:
        total_users = database.get_users_count()
        total_quizzes = len(database.get_all_quizzes())
        total_results = database.get_results_count()
        
        stat_text = (
            "📊 **BOT STATISTIKASI:**\n\n"
            f"👥 **Jami foydalanuvchilar:** {total_users} ta\n"
            f"📂 **Jami yaratilgan testlar:** {total_quizzes} ta\n"
            f"🏁 **Jami yechilgan testlar:** {total_results} marta\n"
        )
        bot.reply_to(message, stat_text, parse_mode='Markdown')
    except Exception as e:
        bot.reply_to(message, f"Statistikani yuklashda xatolik: {e}")

@bot.message_handler(commands=['add_questions'])
def add_more_questions(message):
    if is_group_quiz_active_and_intercept(message):
        return
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    lang = get_user_language(user_id)
    quiz_id = get_current_quiz_id(user_id)
    
    if quiz_id:
        database.set_user_state(user_id, 'waiting_for_questions')
        
        # Get title and count
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        cursor.execute("SELECT title FROM quizzes WHERE quiz_id = ?", (quiz_id,))
        quiz_title = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM questions WHERE quiz_id = ?", (quiz_id,))
        count = cursor.fetchone()[0]
        conn.close()
        
        prompt = T['add_questions_prompt'].get(lang, T['add_questions_prompt']['uz']).format(title=quiz_title, count=count)
        bot.reply_to(message, prompt, reply_markup=get_question_keyboard(lang))
    else:
        bot.reply_to(message, T['no_active_quiz'].get(lang, T['no_active_quiz']['uz']))

@bot.message_handler(commands=['stop'])
def stop_process(message):
    chat_id = message.chat.id
    user_id = message.from_user.id
    lang = get_user_language(user_id)
    
    # 1. Guruh testini to'xtatish
    if message.chat.type in ['group', 'supergroup']:
        if chat_id in group_sessions:
            session = group_sessions[chat_id]
            starter_id = session.get('starter_id')
            if is_user_admin_or_starter(chat_id, user_id, starter_id):
                if 'timer' in session and session['timer']:
                    session['timer'].cancel()
                del group_sessions[chat_id]
                bot.send_message(chat_id, T['stop_group_quiz'].get(lang, T['stop_group_quiz']['uz']), reply_markup=ReplyKeyboardRemove())
            else:
                bot.reply_to(message, T['stop_group_unauthorized'].get(lang, T['stop_group_unauthorized']['uz']))
        else:
            bot.send_message(chat_id, T['no_active_process'].get(lang, T['no_active_process']['uz']))
        return

    # 2. Yakkaxon testni to'xtatish
    if user_id in taking_quiz_sessions:
        session = taking_quiz_sessions[user_id]
        if 'timer' in session and session['timer']:
            session['timer'].cancel()
        del taking_quiz_sessions[user_id]
        bot.send_message(chat_id, T['stop_solo_quiz'].get(lang, T['stop_solo_quiz']['uz']), reply_markup=ReplyKeyboardRemove())
        return

    # 3. Test yaratishni to'xtatish
    state = database.get_user_state(user_id)
    if state != 'none':
        database.set_user_state(user_id, 'none')
        bot.send_message(chat_id, T['stop_creation'].get(lang, T['stop_creation']['uz']), reply_markup=ReplyKeyboardRemove())
        return

    bot.send_message(chat_id, T['no_active_process'].get(lang, T['no_active_process']['uz']))

@bot.message_handler(commands=['done'])
def finish_quiz_cmd(message):
    if is_group_quiz_active_and_intercept(message):
        return
    if message.chat.type in ['group', 'supergroup']: return
    ask_for_time_limit(message)

finish_buttons = {'Tugatish', 'Finish', 'Завершить', 'Bitir', 'Beenden', 'إنهاء'}
@bot.message_handler(func=lambda message: message.text in finish_buttons)
def finish_quiz_btn(message):
    if message.chat.type in ['group', 'supergroup']: return
    ask_for_time_limit(message)

def ask_for_time_limit(message):
    user_id = message.from_user.id
    state = database.get_user_state(user_id)
    
    if state == 'waiting_for_questions':
        database.set_user_state(user_id, 'waiting_for_time_limit')
        
        markup = InlineKeyboardMarkup()
        markup.row(InlineKeyboardButton("10 s", callback_data="set_time_10"),
                   InlineKeyboardButton("15 s", callback_data="set_time_15"),
                   InlineKeyboardButton("20 s", callback_data="set_time_20"))
        markup.row(InlineKeyboardButton("25 s", callback_data="set_time_25"),
                   InlineKeyboardButton("30 s", callback_data="set_time_30"))
        markup.row(InlineKeyboardButton("50 s", callback_data="set_time_50"),
                   InlineKeyboardButton("120 s", callback_data="set_time_120"))
                   
        bot.send_message(message.chat.id, "Bosh menyudasiz.", reply_markup=ReplyKeyboardRemove())
        bot.reply_to(message, "Iltimos, har bir savol uchun qancha vaqt berilishini tanlang:", reply_markup=markup)
    else:
        bot.reply_to(message, "Siz hozir hech qanday test yaratmayapsiz. Boshlash uchun /newquiz ni bosing.", reply_markup=ReplyKeyboardRemove())

@bot.callback_query_handler(func=lambda call: call.data.startswith("set_time_"))
def set_quiz_time_limit(call):
    user_id = call.from_user.id
    state = database.get_user_state(user_id)
    
    if state == 'waiting_for_time_limit':
        time_limit = int(call.data.split('_')[2])
        quiz_id = get_current_quiz_id(user_id)
        
        if quiz_id:
            conn = sqlite3.connect(database.DB_PATH)
            cursor = conn.cursor()
            cursor.execute("UPDATE questions SET time_limit = ? WHERE quiz_id = ?", (time_limit, quiz_id))
            conn.commit()
            conn.close()
            
            database.set_user_state(user_id, 'waiting_for_shuffle')
            
            markup = InlineKeyboardMarkup()
            markup.add(InlineKeyboardButton("Barchasi", callback_data="set_shuffle_3"))
            markup.add(InlineKeyboardButton("Faqat savollar", callback_data="set_shuffle_1"))
            markup.add(InlineKeyboardButton("Faqat variantlar", callback_data="set_shuffle_2"))
            markup.add(InlineKeyboardButton("Yo'q", callback_data="set_shuffle_0"))
            
            bot.edit_message_text(f"⏳ Vaqt {time_limit} soniya qilib belgilandi.\n\nSavollar va variantlar aralashtirilsinmi?", 
                                  chat_id=call.message.chat.id, message_id=call.message.message_id, reply_markup=markup)
        else:
            bot.answer_callback_query(call.id, "Xatolik yuz berdi. (Test topilmadi)")
    else:
        bot.answer_callback_query(call.id, "Bu tugma hozir ishlamaydi.")

@bot.callback_query_handler(func=lambda call: call.data.startswith("set_shuffle_"))
def set_quiz_shuffle(call):
    user_id = call.from_user.id
    state = database.get_user_state(user_id)
    
    if state == 'waiting_for_shuffle':
        shuffle_mode = int(call.data.split('_')[2])
        quiz_id = get_current_quiz_id(user_id)
        
        if quiz_id:
            conn = sqlite3.connect(database.DB_PATH)
            cursor = conn.cursor()
            cursor.execute("UPDATE quizzes SET shuffle_mode = ? WHERE quiz_id = ?", (shuffle_mode, quiz_id))
            
            # Fetch the total number of questions and quiz title
            cursor.execute("SELECT title FROM quizzes WHERE quiz_id = ?", (quiz_id,))
            quiz_title = cursor.fetchone()[0]
            cursor.execute("SELECT COUNT(*) FROM questions WHERE quiz_id = ?", (quiz_id,))
            count = cursor.fetchone()[0]
            
            conn.commit()
            conn.close()
            
            database.set_user_state(user_id, 'none')
            
            markup = InlineKeyboardMarkup()
            start_url = f"https://t.me/{bot_info.username}?start=quiz_{quiz_id}"
            group_url = f"https://t.me/{bot_info.username}?startgroup=quiz_{quiz_id}"
            
            markup.add(InlineKeyboardButton("Yakkaxon yechish 👤", url=start_url))
            markup.add(InlineKeyboardButton("Guruhda yechish 👥", url=group_url))
            markup.add(InlineKeyboardButton("Do'stlarga ulashish ↗️", url=f"https://t.me/share/url?url={start_url}"))
            
            success_text = (
                f"✅ Test muvaffaqiyatli yaratildi va saqlandi!\n\n"
                f"Sizning ushbu \"{quiz_title}\" testingizda {count} ta test (savol) bo'ldi. 🎉\n\n"
                f"Quyidagi tugmalar orqali testni o'zingiz yechishingiz yoki guruhlarga yuborishingiz mumkin:"
            )
            
            bot.edit_message_text(success_text, 
                                  chat_id=call.message.chat.id, message_id=call.message.message_id, reply_markup=markup)
            bot.answer_callback_query(call.id, "Aralashtirish rejimi saqlandi!")
        else:
            bot.answer_callback_query(call.id, "Xatolik yuz berdi.")
    else:
        bot.answer_callback_query(call.id, "Bu tugma hozir ishlamaydi.")

@bot.message_handler(commands=['language'])
def command_change_language(message):
    if is_group_quiz_active_and_intercept(message):
        return
    user_id = message.from_user.id
    lang = get_user_language(user_id)
    bot.reply_to(message, T['choose_language'].get(lang, T['choose_language']['uz']), reply_markup=get_language_keyboard())

@bot.message_handler(commands=['skip'])
def skip_description(message):
    if is_group_quiz_active_and_intercept(message):
        return
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    state = database.get_user_state(user_id)
    lang = get_user_language(user_id)
    
    if state == 'waiting_for_description':
        database.set_user_state(user_id, 'waiting_for_questions')
        text = T['description_skipped'].get(lang, T['description_skipped']['uz']) + T['ask_for_questions_prompt'].get(lang, T['ask_for_questions_prompt']['uz'])
        bot.reply_to(message, text, reply_markup=get_question_keyboard(lang))

def get_button_type(text):
    labels = {
        'uz': ('➕ Yangi test yaratish', '📂 Mening testlarim', '📊 Natijalar', '🌐 Tilni o\'zgartirish'),
        'en': ('➕ Create New Quiz', '📂 My Quizzes', '📊 Results', '🌐 Change Language'),
        'ru': ('➕ Создать новый тест', '📂 Мои тесты', '📊 Результаты', '🌐 Сменить язык'),
        'tr': ('➕ Yeni Test Oluştur', '📂 Testlerim', '📊 Sonuçlar', '🌐 Dili Değiştir'),
        'de': ('➕ Neuen Test erstellen', '📂 Meine Tests', '📊 Ergebnisse', '🌐 Sprache ändern'),
        'ar': ('➕ إنشاء اختبار جديد', '📂 اختباراتي', '📊 النتائج', '🌐 تغيير اللغة')
    }
    for lang, (new_quiz, my_quizzes, results, change_lang) in labels.items():
        if text == new_quiz: return 'new_quiz'
        if text == my_quizzes: return 'my_quizzes'
        if text == results: return 'results'
        if text == change_lang: return 'change_lang'
    return None

@bot.message_handler(content_types=['text'])
def handle_text(message):
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    state = database.get_user_state(user_id)
    lang = get_user_language(user_id)
    
    if state == 'waiting_for_title':
        title = message.text
        quiz_id = database.create_quiz(user_id, title)
        user_sessions[user_id] = {'quiz_id': quiz_id}
        database.set_user_state(user_id, 'waiting_for_description')
        bot.reply_to(message, T['waiting_for_description_prompt'].get(lang, T['waiting_for_description_prompt']['uz']))
        
    elif state == 'waiting_for_description':
        description = message.text
        quiz_id = get_current_quiz_id(user_id)
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        cursor.execute("UPDATE quizzes SET description = ? WHERE quiz_id = ?", (description, quiz_id))
        conn.commit()
        conn.close()
        database.set_user_state(user_id, 'waiting_for_questions')
        text = T['description_saved'].get(lang, T['description_saved']['uz']) + T['ask_for_questions_prompt'].get(lang, T['ask_for_questions_prompt']['uz'])
        bot.reply_to(message, text, reply_markup=get_question_keyboard(lang))
        
    else:
        button_type = get_button_type(message.text)
        if button_type == 'new_quiz':
            create_new_quiz(message)
        elif button_type == 'my_quizzes':
            show_my_quizzes(message)
        elif button_type == 'results':
            bot.reply_to(message, T['results_section_coming_soon'].get(lang, T['results_section_coming_soon']['uz']))
        elif button_type == 'change_lang':
            bot.reply_to(message, T['choose_language'].get(lang, T['choose_language']['uz']), reply_markup=get_language_keyboard())
        else:
            bot.reply_to(message, T['main_menu'].get(lang, T['main_menu']['uz']), reply_markup=get_main_keyboard(lang))

def show_my_quizzes(message):
    user_id = message.from_user.id
    quizzes = database.get_my_quizzes(user_id)
    lang = get_user_language(user_id)
    
    if not quizzes:
        bot.reply_to(message, T['no_quizzes_yet'].get(lang, T['no_quizzes_yet']['uz']))
        return
        
    text = T['my_quizzes_title'].get(lang, T['my_quizzes_title']['uz'])
    markup = InlineKeyboardMarkup()
    solve_btn_label = T['solve_quiz'].get(lang, T['solve_quiz']['uz'])
    for q_id, title in quizzes:
        text += f"🔹 {title}\n"
        start_url = f"https://t.me/{bot_info.username}?start=quiz_{q_id}"
        markup.add(InlineKeyboardButton(f"{solve_btn_label}: {title}", url=start_url))
        
    bot.reply_to(message, text, reply_markup=markup, parse_mode='Markdown')



@bot.message_handler(content_types=['poll'])
def handle_poll(message):
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    state = database.get_user_state(user_id)
    
    if state == 'waiting_for_questions':
        poll = message.poll
        if not poll.type == 'quiz':
            bot.reply_to(message, "Iltimos, so'rovnomani **Quiz Mode** (Viktorina rejimi) da yarating!")
            return
            
        quiz_id = get_current_quiz_id(user_id)
        question_text = poll.question
        options = json.dumps([option.text for option in poll.options])
        correct_option_id = poll.correct_option_id
        explanation = poll.explanation if poll.explanation else ""
        
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        cursor.execute("""
            INSERT INTO questions (quiz_id, question_text, options, correct_option_id, explanation)
            VALUES (?, ?, ?, ?, ?)
        """, (quiz_id, question_text, options, correct_option_id, explanation))
        conn.commit()
        
        # Get title and count
        cursor.execute("SELECT title FROM quizzes WHERE quiz_id = ?", (quiz_id,))
        quiz_title = cursor.fetchone()[0]
        cursor.execute("SELECT COUNT(*) FROM questions WHERE quiz_id = ?", (quiz_id,))
        count = cursor.fetchone()[0]
        conn.close()
        
        response_text = (
            f"Yaxshi. \"{quiz_title}\" testingizda hozirda {count} ta savol bor. "
            "Agar oxirgi savolda xatoga yo'l qo'ygan bo'lsangiz, /undo buyrug'ini yuborish orqali orqaga qaytarishingiz mumkin.\n\n"
            "Keyingi savolni yaratish uchun yana **«Savol yaratish»** tugmasini bosing yoki tugatish uchun **«Tugatish»** tugmasini bosing."
        )
        bot.reply_to(message, response_text, reply_markup=get_question_keyboard())

@bot.message_handler(commands=['undo'])
def undo_last_question(message):
    if is_group_quiz_active_and_intercept(message):
        return
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    state = database.get_user_state(user_id)
    
    if state != 'waiting_for_questions':
        bot.reply_to(message, "Siz hozir savol qo'shish jarayonida emassiz.")
        return
        
    quiz_id = get_current_quiz_id(user_id)
    if not quiz_id:
        bot.reply_to(message, "Yaratilayotgan test topilmadi.")
        return
        
    conn = sqlite3.connect(database.DB_PATH)
    cursor = conn.cursor()
    
    # Get quiz title
    cursor.execute("SELECT title FROM quizzes WHERE quiz_id = ?", (quiz_id,))
    quiz_res = cursor.fetchone()
    if not quiz_res:
        conn.close()
        bot.reply_to(message, "Test topilmadi.")
        return
    quiz_title = quiz_res[0]
    
    # Get the last question
    cursor.execute("SELECT question_id, question_text FROM questions WHERE quiz_id = ? ORDER BY question_id DESC LIMIT 1", (quiz_id,))
    q_res = cursor.fetchone()
    
    if not q_res:
        conn.close()
        bot.reply_to(message, f"\"{quiz_title}\" testingizda hali hech qanday savol yo'q, shuning uchun /undo qilib bo'lmaydi.")
        return
        
    q_id, q_text = q_res
    cursor.execute("DELETE FROM questions WHERE question_id = ?", (q_id,))
    conn.commit()
    
    # Get the new count
    cursor.execute("SELECT COUNT(*) FROM questions WHERE quiz_id = ?", (quiz_id,))
    new_count = cursor.fetchone()[0]
    conn.close()
    
    response_text = (
        f"🗑 Oxirgi savol muvaffaqiyatli o'chirildi!\n"
        f"O'chirilgan savol: \"{q_text}\"\n\n"
        f"Hozirda \"{quiz_title}\" testingizda {new_count} ta savol qoldi.\n"
        f"Yangi savol yaratish uchun pastdagi **«Savol yaratish»** tugmasini bosing."
    )
    bot.reply_to(message, response_text, reply_markup=get_question_keyboard())

# --- AI ORQALI TEST YARATISH (Fayl va Rasmlar) ---

@bot.message_handler(content_types=['document'])
def handle_document(message):
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    
    ADMIN_ID = 6559589296
    
    # Limit tekshirish
    if user_id != ADMIN_ID:
        usage = database.get_daily_ai_usage(user_id)
        if usage >= 3:
            bot.reply_to(message, "❌ Sizning bugungi limitingiz tugadi. Ertaga fayl yoki rasmlaringizni qaytadan yuborib ko'ring.")
            return
            
    # API key mavjudligini tekshirish
    if not os.getenv('GEMINI_API_KEY'):
        text = (
            "⚠️ **AI Quiz funksiyasi faol emas!**\n\n"
            "Ushbu funksiyadan foydalanish uchun bot egasi `.env` fayliga `GEMINI_API_KEY` kalitini qo'shishi kerak.\n\n"
            "🔑 **Kalitni qanday olish mumkin?**\n"
            "1. https://aistudio.google.com/ saytiga kiring.\n"
            "2. Bepul **Gemini API Key** oling.\n"
            "3. Botning `.env` fayliga yozing:\n"
            "`GEMINI_API_KEY=sizning_kalitingiz`"
        )
        bot.reply_to(message, text, parse_mode='Markdown')
        return

    file_name = message.document.file_name
    file_ext = os.path.splitext(file_name)[1].lower()
    
    if file_ext not in ['.pdf', '.docx', '.txt', '.pptx']:
        bot.reply_to(message, "❌ Kechirasiz, faqat `.pdf`, `.docx`, `.pptx` va `.txt` formatidagi fayllarni qabul qila olaman.")
        return
        
    processing_msg = bot.reply_to(message, "📥 **Hujjat yuklab olinmoqda va o'qilmoqda...** Iltimos, kuting ⏳", parse_mode='Markdown')
    
    try:
        file_info = bot.get_file(message.document.file_id)
        downloaded_file = bot.download_file(file_info.file_path)
        
        # Faylni vaqtinchalik saqlash
        temp_path = f"temp_{user_id}{file_ext}"
        with open(temp_path, 'wb') as new_file:
            new_file.write(downloaded_file)
            
        # Matnni ajratib olish
        content_text = ""
        if file_ext == '.pdf':
            content_text = extract_text_from_pdf(temp_path)
        elif file_ext == '.docx':
            content_text = extract_text_from_docx(temp_path)
        elif file_ext == '.pptx':
            content_text = extract_text_from_pptx(temp_path)
        elif file_ext == '.txt':
            content_text = downloaded_file.decode('utf-8', errors='ignore')
            
        # Vaqtinchalik faylni o'chirish
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        if not content_text or len(content_text.strip()) < 20:
            bot.edit_message_text("❌ Fayldan matn o'qib bo'lmadi yoki undagi matn juda qisqa (kamida 20 ta belgi bo'lishi kerak).", chat_id=message.chat.id, message_id=processing_msg.message_id)
            return
            
        # Sessiyaga yozish
        ai_sessions[user_id] = {
            'type': 'text',
            'data': content_text,
            'file_name': file_name
        }
        
        # Savollar sonini tanlash tugmalari
        markup = InlineKeyboardMarkup()
        markup.row(InlineKeyboardButton("5 ta savol 📝", callback_data="ai_count_5"),
                   InlineKeyboardButton("10 ta savol 📝", callback_data="ai_count_10"))
        markup.row(InlineKeyboardButton("15 ta savol 📝", callback_data="ai_count_15"),
                   InlineKeyboardButton("20 ta savol 📝", callback_data="ai_count_20"))
        markup.row(InlineKeyboardButton("30 ta savol 📝", callback_data="ai_count_30"),
                   InlineKeyboardButton("40 ta savol 📝", callback_data="ai_count_40"))
        
        limit_text = ""
        if user_id != ADMIN_ID:
            usage = database.get_daily_ai_usage(user_id)
            limit_text = f"\n⚠️ *Eslatma: Siz bir kunda maksimal 3 tagacha fayl yoki rasm yubora olasiz. (Bugungi qolgan limitingiz: {3 - usage} ta)*\n"
        else:
            limit_text = f"\n⭐ *Siz uchun AI limiti cheksiz!*\n"
            
        success_text = (
            f"📄 **Fayl muvaffaqiyatli o'qildi!**\n"
            f"Hujjat: `{file_name}`\n"
            f"Matn hajmi: {len(content_text)} ta belgi.\n"
            f"{limit_text}\n"
            f"Ushbu fayl asosida nechta test savoli yaratmoqchisiz? Tanlang 👇"
        )
        bot.edit_message_text(success_text, chat_id=message.chat.id, message_id=processing_msg.message_id, reply_markup=markup, parse_mode='Markdown')
        
    except Exception as e:
        bot.edit_message_text(f"❌ Faylni yuklash yoki o'qishda xatolik yuz berdi: {e}", chat_id=message.chat.id, message_id=processing_msg.message_id)


@bot.message_handler(content_types=['photo'])
def handle_photo(message):
    if message.chat.type in ['group', 'supergroup']: return
    user_id = message.from_user.id
    
    ADMIN_ID = 6559589296
    
    # Limit tekshirish
    if user_id != ADMIN_ID:
        usage = database.get_daily_ai_usage(user_id)
        if usage >= 3:
            bot.reply_to(message, "❌ Sizning bugungi limitingiz tugadi. Ertaga fayl yoki rasmlaringizni qaytadan yuborib ko'ring.")
            return
            
    # API key mavjudligini tekshirish
    if not os.getenv('GEMINI_API_KEY'):
        text = (
            "⚠️ **AI Quiz funksiyasi faol emas!**\n\n"
            "Ushbu funksiyadan foydalanish uchun bot egasi `.env` fayliga `GEMINI_API_KEY` kalitini qo'shishi kerak.\n\n"
            "🔑 **Kalitni qanday olish mumkin?**\n"
            "1. https://aistudio.google.com/ saytiga kiring.\n"
            "2. Bepul **Gemini API Key** oling.\n"
            "3. Botning `.env` fayliga yozing:\n"
            "`GEMINI_API_KEY=sizning_kalitingiz`"
        )
        bot.reply_to(message, text, parse_mode='Markdown')
        return

    processing_msg = bot.reply_to(message, "📥 **Rasm yuklab olinmoqda...** Iltimos, kuting ⏳", parse_mode='Markdown')
    
    try:
        # Eng katta rasmni olish
        photo = message.photo[-1]
        file_info = bot.get_file(photo.file_id)
        downloaded_file = bot.download_file(file_info.file_path)
        
        # Sessiyaga yozish (Gemini uchun bayt ko'rinishida)
        ai_sessions[user_id] = {
            'type': 'image',
            'data': downloaded_file,
            'file_name': 'Rasm'
        }
        
        # Savollar sonini tanlash tugmalari
        markup = InlineKeyboardMarkup()
        markup.row(InlineKeyboardButton("5 ta savol 📝", callback_data="ai_count_5"),
                   InlineKeyboardButton("10 ta savol 📝", callback_data="ai_count_10"))
        markup.row(InlineKeyboardButton("15 ta savol 📝", callback_data="ai_count_15"),
                   InlineKeyboardButton("20 ta savol 📝", callback_data="ai_count_20"))
        markup.row(InlineKeyboardButton("30 ta savol 📝", callback_data="ai_count_30"),
                   InlineKeyboardButton("40 ta savol 📝", callback_data="ai_count_40"))
        
        limit_text = ""
        if user_id != ADMIN_ID:
            usage = database.get_daily_ai_usage(user_id)
            limit_text = f"\n⚠️ *Eslatma: Siz bir kunda maksimal 3 tagacha fayl yoki rasm yubora olasiz. (Bugungi qolgan limitingiz: {3 - usage} ta)*\n"
        else:
            limit_text = f"\n⭐ *Siz uchun AI limiti cheksiz!*\n"
            
        success_text = (
            f"🖼 **Rasm muvaffaqiyatli qabul qilindi!**\n"
            f"{limit_text}\n"
            f"Ushbu rasm/skrinshotdagi ma'lumotlar asosida nechta test savoli yaratmoqchisiz? Tanlang 👇"
        )
        bot.edit_message_text(success_text, chat_id=message.chat.id, message_id=processing_msg.message_id, reply_markup=markup, parse_mode='Markdown')
        
    except Exception as e:
        bot.edit_message_text(f"❌ Rasmni yuklashda xatolik yuz berdi: {e}", chat_id=message.chat.id, message_id=processing_msg.message_id)


@bot.callback_query_handler(func=lambda call: call.data.startswith("ai_count_"))
def handle_ai_count(call):
    user_id = call.from_user.id
    session = ai_sessions.get(user_id)
    
    if not session:
        bot.answer_callback_query(call.id, "Sessiya eskirgan. Qaytadan urinib ko'ring.")
        bot.edit_message_text("❌ Hech qanday faol AI sessiyasi topilmadi. Qayta urinib ko'ring (botga yangi fayl yoki rasm yuboring).", 
                              chat_id=call.message.chat.id, message_id=call.message.message_id)
        return
        
    count = int(call.data.split('_')[2])
    session['count'] = count
    bot.answer_callback_query(call.id, f"{count} ta savol tanlandi!")
    
    # Vaqtni tanlash tugmalari
    markup = InlineKeyboardMarkup()
    markup.row(InlineKeyboardButton("10 s ⏱️", callback_data="ai_time_10"),
               InlineKeyboardButton("15 s ⏱️", callback_data="ai_time_15"),
               InlineKeyboardButton("20 s ⏱️", callback_data="ai_time_20"))
    markup.row(InlineKeyboardButton("25 s ⏱️", callback_data="ai_time_25"),
               InlineKeyboardButton("50 s ⏱️", callback_data="ai_time_50"),
               InlineKeyboardButton("120 s ⏱️", callback_data="ai_time_120"))
               
    text = (
        f"🎯 Savollar soni: **{count} ta** tanlandi.\n\n"
        f"Endi, har bir savol uchun qancha vaqt berilishini tanlang 👇"
    )
    bot.edit_message_text(text, chat_id=call.message.chat.id, message_id=call.message.message_id, reply_markup=markup, parse_mode='Markdown')


@bot.callback_query_handler(func=lambda call: call.data.startswith("ai_time_"))
def handle_ai_time(call):
    user_id = call.from_user.id
    session = ai_sessions.get(user_id)
    
    if not session or 'count' not in session:
        bot.answer_callback_query(call.id, "Sessiya eskirgan. Qaytadan urinib ko'ring.")
        bot.edit_message_text("❌ Hech qanday faol AI sessiyasi topilmadi. Qayta urinib ko'ring.", 
                              chat_id=call.message.chat.id, message_id=call.message.message_id)
        return
        
    time_limit = int(call.data.split('_')[2])
    count = session['count']
    session['time_limit'] = time_limit
    
    bot.answer_callback_query(call.id, f"{time_limit} soniya tanlandi!")
    
    bot.edit_message_text(f"🤖 **AI ma'lumotlarni tahlil qilmoqda va test yaratmoqda...**\n"
                          f"Savollar soni: {count} ta\n"
                          f"Vaqt limiti: {time_limit} soniya\n\n"
                          f"Iltimos, kuting. Bu jarayon 10-30 soniya vaqt olishi mumkin ⏳", 
                          chat_id=call.message.chat.id, message_id=call.message.message_id, parse_mode='Markdown')
    
    threading.Thread(target=generate_ai_quiz_thread, args=(call.message, user_id, session, count, time_limit)).start()


def generate_ai_quiz_thread(message, user_id, session, count, time_limit):
    chat_id = message.chat.id
    message_id = message.message_id
    
    try:
        model = genai.GenerativeModel("gemini-2.5-flash")
        
        prompt = (
            f"Sizga taqdim etilgan matn yoki rasmdagi ma'lumotlardan foydalanib, roppa-rosa {count} ta "
            "ko'p variantli (multiple choice) test savollarini yarating. "
            "Barcha savollar, variantlar va tushuntirishlar faqat o'zbek tilida bo'lishi kerak. "
            "Har bir savol uchun aniq 4 ta variant (options) tayyorlang.\n\n"
            "Agar matnda tayyor savollar va to'g'ri javoblar ko'rsatilgan bo'lsa (masalan `#`, `*`, `bold` belgilar orqali), "
            "o'sha tayyor savollardan va ko'rsatilgan to'g'ri javoblardan foydalaning va to'g'ri javobning indeksini belgilang.\n\n"
            "Natijani mutlaqo qat'iy ravishda quyidagi JSON formatida qaytaring. Hech qanday boshqa matn, kirish yoki markdown formatting (masalan ```json) qo'shmang, "
            "faqat quyidagi struktura bo'yicha toza JSON yuboring:\n"
            "{\n"
            '  "title": "Mavzuga mos chiroyli test sarlavhasi (maksimal 50 ta belgi)",\n'
            '  "description": "Fayl/mavzu haqida qisqacha izoh (maksimal 150 ta belgi)",\n'
            '  "questions": [\n'
            "    {\n"
            '      "question": "Savol matni (maksimal 150 ta belgi)",\n'
            '      "options": ["1-variant", "2-variant", "3-variant", "4-variant"],\n'
            '      "correct_option_index": 0,\n'
            '      "explanation": "Nega aynan shu variant to\'g\'ri ekanligining qisqacha izohi (maksimal 100 ta belgi, ixtiyoriy)"\n'
            "    }\n"
            "  ]\n"
            "}"
        )
        
        # Gemini modeliga so'rov yuborish
        if session['type'] == 'text':
            full_content = f"Hujjat nomi: {session['file_name']}\n\nHujjat matni:\n{session['data']}"
            if len(full_content) > 60000:
                full_content = full_content[:60000] + "\n...[Matn juda uzunligi uchun kesildi]..."
                
            response = model.generate_content(
                [prompt, full_content],
                generation_config={"response_mime_type": "application/json"}
            )
        else:
            image_part = {
                "mime_type": "image/jpeg",
                "data": session['data']
            }
            response = model.generate_content(
                [prompt, image_part],
                generation_config={"response_mime_type": "application/json"}
            )
            
        raw_text = response.text.strip()
        
        # Agar rasm/matn mos kelmasa va xatolik bo'lsa
        if not raw_text:
            bot.edit_message_text("❌ Xatolik: AI matndan ma'lumot ololmadi yoki javob qaytarmadi.", chat_id=chat_id, message_id=message_id)
            return
            
        # Kutilmagan markdown belgilarni tozalash (agar response_mime_type ishlamasa)
        if raw_text.startswith("```json"):
            raw_text = raw_text.replace("```json", "", 1)
        if raw_text.endswith("```"):
            raw_text = raw_text[:-3].strip()
        raw_text = raw_text.strip()
        
        data = json.loads(raw_text)
        title = data.get('title', f"AI Test - {session['file_name']}")
        description = data.get('description', "AI tomonidan avtomatik yaratilgan test.")
        questions = data.get('questions', [])
        
        if not questions:
            bot.edit_message_text("❌ Xatolik: AI birorta ham test savoli yarata olmadi. Boshqa fayl yoki rasm yuborib ko'ring.", chat_id=chat_id, message_id=message_id)
            return
            
        # Testni bazada yaratish
        quiz_id = database.create_quiz(user_id, title)
        
        conn = sqlite3.connect(database.DB_PATH)
        cursor = conn.cursor()
        cursor.execute("UPDATE quizzes SET description = ? WHERE quiz_id = ?", (description, quiz_id))
        
        added_count = 0
        for q in questions:
            q_text = q.get('question')
            options = q.get('options')
            correct_idx = q.get('correct_option_index', 0)
            explanation = q.get('explanation', '')
            
            if q_text and options and len(options) >= 2:
                options_json = json.dumps(options[:10]) # Telegram max 10 ta variant qabul qiladi
                cursor.execute("""
                    INSERT INTO questions (quiz_id, question_text, options, correct_option_id, explanation, time_limit)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (quiz_id, q_text, options_json, correct_idx, explanation, time_limit))
                added_count += 1
                
        conn.commit()
        conn.close()
        
        # Limit hisoblagichini oshirish
        database.increment_daily_ai_usage(user_id)
        
        # Sessiyani tozalash
        if user_id in ai_sessions:
            del ai_sessions[user_id]
            
        if added_count == 0:
            bot.edit_message_text("❌ Xatolik: Savollarni tahlil qilishda yoki bazaga yozishda muammo yuz berdi.", chat_id=chat_id, message_id=message_id)
            return
            
        # O'ynash havolalarini generatsiya qilish
        start_url = f"https://t.me/{bot_info.username}?start=quiz_{quiz_id}"
        group_url = f"https://t.me/{bot_info.username}?startgroup=quiz_{quiz_id}"
        
        markup = InlineKeyboardMarkup()
        markup.add(InlineKeyboardButton("Yakkaxon yechish 👤", url=start_url))
        markup.add(InlineKeyboardButton("Guruhda yechish 👥", url=group_url))
        markup.add(InlineKeyboardButton("Do'stlarga ulashish ↗️", url=f"https://t.me/share/url?url={start_url}"))
        
        success_text = (
            f"🎉 **AI Quiz muvaffaqiyatli yaratildi!**\n\n"
            f"📌 **Sarlavha:** {title}\n"
            f"📝 **Izoh:** {description}\n"
            f"❓ **Savollar soni:** {added_count} ta\n\n"
            f"Quyidagi tugmalar orqali testni yechishingiz yoki ulashishingiz mumkin:"
        )
        bot.edit_message_text(success_text, chat_id=chat_id, message_id=message_id, reply_markup=markup, parse_mode='Markdown')
        
    except json.JSONDecodeError:
        bot.edit_message_text("❌ Xatolik: AI qaytargan ma'lumotni JSON formatida o'qib bo'lmadi. Qayta urinib ko'ring.", chat_id=chat_id, message_id=message_id)
    except Exception as e:
        error_msg = str(e)
        print(f"AI error: {error_msg}")
        if "429" in error_msg or "Quota" in error_msg:
            user_msg = "❌ Kechirasiz, ayni vaqtda botga so'rovlar juda ko'payib ketdi (AI limiti tugadi). Iltimos, birozdan so'ng (1-2 daqiqa) qayta urinib ko'ring."
        else:
            user_msg = "❌ Test yaratishda kutilmagan xatolik yuz berdi. Iltimos, qaytadan urinib ko'ring."
        bot.edit_message_text(user_msg, chat_id=chat_id, message_id=message_id)



import time
print("Bot ishga tushmoqda (Versiya: 1.1 - FixID qo'shildi)...")
print(f"Baza yo'li: {database.DB_PATH}")
while True:
    try:
        bot.polling(none_stop=True, timeout=60, long_polling_timeout=60)
    except Exception as e:
        print(f"Internet uzilishi yoki xatolik: {e}")
        time.sleep(3)
